# SECURITY WARNING:
# The file tools below accept absolute and relative paths and searches files broadly in 
# common folders like 
# (Downloads, Documents, C:/Users/yourusername, C:/Users/yourusername/anysubfolder, Videos, Pictures, Pictures/Screeshots, 
# etc) if paths are not provided in the user prompts.
# This allows reading, writing, deleting, and renaming files anywhere on the system
# that the process has permission for. Use with caution!

import os
import threading
import requests
import json
import base64
import webbrowser
import imaplib
import time
import re
import subprocess
import platform
import uuid
import cv2
import shutil
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
from email.utils import formatdate, make_msgid
from email.policy import SMTP
from urllib.parse import quote
from openai import OpenAI
from typing import Dict, List, Callable, Any, Optional, TypedDict
from django.conf import settings
from datetime import datetime, timedelta
from decimal import Decimal

from langchain_openai import ChatOpenAI
from langchain_core.tools import StructuredTool
from langchain_core.messages import (
    BaseMessage, HumanMessage, AIMessage, ToolMessage, SystemMessage
)
from langgraph.graph import StateGraph, END
from pydantic import create_model, Field

from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, HRFlowable
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors as rl_colors
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_LEFT, TA_CENTER
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from openpyxl import Workbook
from openpyxl.styles import Font as XlFont, PatternFill, Alignment as XlAlignment, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN

class AgentControlState:
    def __init__(self):
        self.stopped = False
        self.tools_enabled = True
        self.lock = threading.Lock()

    def stop(self):
        with self.lock:
            self.stopped = True

    def disable_tools(self):
        with self.lock:
            self.tools_enabled = False

    def enable_tools(self):
        with self.lock:
            self.tools_enabled = True


# ─── Duffel REST API Helpers ──────────────────────────────────────────

DUFFEL_API_BASE = "https://api.duffel.com"

# Caches for cross-tool state (search -> book)
_duffel_offer_cache = {}   # offer_id -> {"passenger_ids": [...], "total_amount": str, "total_currency": str, "expires_at": str}


def _duffel_headers():
    """Return authorization headers for all Duffel REST API calls."""
    token = os.getenv('DUFFEL_API_TOKEN', '')
    if not token:
        try:
            from django.conf import settings as djsettings
            token = getattr(djsettings, 'DUFFEL_API_TOKEN', '')
        except Exception:
            pass
    if not token:
        raise ValueError("DUFFEL_API_TOKEN is not set in the environment.")
    return {
        "Authorization": f"Bearer {token}",
        "Duffel-Version": "v2",
        "Accept": "application/json",
        "Content-Type": "application/json",
    }


def _duffel_post(path, payload, timeout=30):
    """POST to the Duffel API and return (status_code, response_data_dict)."""
    resp = requests.post(f"{DUFFEL_API_BASE}{path}", headers=_duffel_headers(), json=payload, timeout=timeout)
    return resp.status_code, resp.json()


def _duffel_get(path, params=None, timeout=15):
    """GET from the Duffel API and return (status_code, response_data_dict)."""
    resp = requests.get(f"{DUFFEL_API_BASE}{path}", headers=_duffel_headers(), params=params, timeout=timeout)
    return resp.status_code, resp.json()


def _parse_iso_duration(iso_dur):
    """Convert ISO 8601 duration like 'PT2H26M' to '2h 26m'."""
    if not iso_dur:
        return ''
    m = re.match(r'PT(?:(\d+)H)?(?:(\d+)M)?', iso_dur)
    if not m:
        return iso_dur
    hours = m.group(1) or '0'
    minutes = m.group(2) or '0'
    return f"{hours}h {minutes}m"


class ToolDefinition:
    def __init__(self, name: str, description: str, parameters: Dict[str, Any], function: Callable[[Dict[str, Any]], str], requires_approval: bool = False):
        self.name = name
        self.description = description
        self.parameters = parameters
        self.function = function
        self.requires_approval = requires_approval


class ApprovalAwareTool(StructuredTool):
    """StructuredTool subclass that preserves the requires_approval flag."""
    requires_approval: bool = False


def _json_type_to_python(json_type: str):
    """Map JSON Schema types to Python types."""
    mapping = {
        "string": str,
        "integer": int,
        "number": float,
        "boolean": bool,
        "array": list,
        "object": dict,
    }
    return mapping.get(json_type, str)


def tool_definition_to_langchain(td: ToolDefinition) -> ApprovalAwareTool:
    """Convert a legacy ToolDefinition into a LangChain ApprovalAwareTool."""
    properties = td.parameters.get("properties", {})
    required_fields = set(td.parameters.get("required", []))

    field_definitions = {}
    for field_name, field_schema in properties.items():
        field_type = _json_type_to_python(field_schema.get("type", "string"))
        description = field_schema.get("description", "")
        if field_name in required_fields:
            field_definitions[field_name] = (field_type, Field(description=description))
        else:
            field_definitions[field_name] = (
                Optional[field_type],
                Field(default=None, description=description),
            )

    ArgsModel = create_model(f"{td.name}_args", **field_definitions)

    original_func = td.function

    def wrapper_func(**kwargs) -> str:
        # Strip None values so the original function only sees provided args
        cleaned = {k: v for k, v in kwargs.items() if v is not None}
        return original_func(cleaned)

    return ApprovalAwareTool(
        name=td.name,
        description=td.description,
        func=wrapper_func,
        args_schema=ArgsModel,
        requires_approval=td.requires_approval,
    )


def find_file_broadly(filename: str) -> str:
    """Attempt to find a file by name in common user directories."""
    if not filename:
        return None
        
    if os.path.isabs(filename):
        return filename if os.path.exists(filename) else None
        
    # Handle nested paths by resolving the base directory first
    normalized_path = filename.replace('\\', '/')
    if '/' in normalized_path:
        parts = normalized_path.split('/')
        base_dir = find_directory_broadly(parts[0])
        if base_dir and os.path.isabs(base_dir):
            return find_file_broadly(os.path.join(base_dir, *parts[1:]))

    # Check for well-known folders directly if filename is one of them
    user_home = os.path.expanduser("~")
    well_known = {
        "frank": user_home,
        "documents": os.path.join(user_home, "Documents"),
        "desktop": os.path.join(user_home, "Desktop"),
        "downloads": os.path.join(user_home, "Downloads"),
        "pictures": os.path.join(user_home, "Pictures"),
        "videos": os.path.join(user_home, "Videos"),
    }
    if filename.lower() in well_known:
        path = well_known[filename.lower()]
        if os.path.exists(path):
            return path

    # Search in current directory
    if os.path.exists(filename):
        return os.path.abspath(filename)
        
    # Define common search paths
    user_home = os.path.expanduser("~")
    search_dirs = [
        user_home,
        os.path.join(user_home, "Documents"),
        os.path.join(user_home, "OneDrive", "Documents"),
        os.path.join(user_home, "Desktop"),
        os.path.join(user_home, "OneDrive", "Desktop"),
        os.path.join(user_home, "Downloads"),
        os.path.join(user_home, "Pictures"),
        os.path.join(user_home, "Pictures", "Screenshots"),
        os.path.join(user_home, "Videos"),
        os.path.join(user_home, "projects"),
        os.path.join(user_home, "repos"),
        os.path.join(user_home, "work"),
        os.getcwd(),
    ]
    
    # Filter out non-existent directories and ensure uniqueness
    search_dirs = list(dict.fromkeys(d for d in search_dirs if os.path.exists(d)))
    
    # Search for exact match first
    for directory in search_dirs:
        # Check direct children
        try:
            for entry in os.listdir(directory):
                if entry.lower() == filename.lower() or os.path.splitext(entry)[0].lower() == filename.lower():
                    full_path = os.path.join(directory, entry)
                    if os.path.isfile(full_path):
                        return full_path
        except Exception:
            continue
            
    # Search recursively with limited depth if not found
    for directory in search_dirs:
        try:
            for root, dirs, files in os.walk(directory):
                # Limit depth to avoid performance issues
                try:
                    rel = os.path.relpath(root, directory)
                    depth = 0 if rel == '.' else len(rel.replace('\\', '/').split('/'))
                except ValueError:
                    depth = 0
                if depth >= 4:
                    del dirs[:] # Don't go deeper
                
                for file in files:
                    if file.lower() == filename.lower() or os.path.splitext(file)[0].lower() == filename.lower() or filename.lower() in file.lower():
                        return os.path.join(root, file)
        except Exception:
            continue
            
    return None
    
    
def find_directory_broadly(dirname: str) -> str:
    """Attempt to find a directory by name in common user directories."""
    if not dirname:
        return None
        
    if os.path.isabs(dirname):
        return dirname if os.path.isdir(dirname) else None
        
    # Handle nested paths by resolving the base directory first
    normalized_path = dirname.replace('\\', '/')
    if '/' in normalized_path:
        parts = normalized_path.split('/')
        base_dir = find_directory_broadly(parts[0])
        if base_dir and os.path.isabs(base_dir):
            return os.path.abspath(os.path.join(base_dir, *parts[1:]))

    # Check for well-known folders directly if dirname is one of them
    user_home = os.path.expanduser("~")
    well_known = {
        "frank": user_home,
        "documents": os.path.join(user_home, "Documents"),
        "desktop": os.path.join(user_home, "Desktop"),
        "downloads": os.path.join(user_home, "Downloads"),
        "pictures": os.path.join(user_home, "Pictures"),
        "videos": os.path.join(user_home, "Videos"),
    }
    if dirname.lower() in well_known:
        path = well_known[dirname.lower()]
        if os.path.exists(path):
            return path

    # Search in current directory
    if os.path.isdir(dirname):
        return os.path.abspath(dirname)
        
    # Define common search paths
    user_home = os.path.expanduser("~")
    search_dirs = [
        os.path.join(user_home, "Documents"),
        os.path.join(user_home, "OneDrive", "Documents"),
        os.path.join(user_home, "Desktop"),
        os.path.join(user_home, "OneDrive", "Desktop"),
        os.path.join(user_home, "Downloads"),
        os.path.join(user_home, "Pictures"),
        os.path.join(user_home, "Videos"),
        os.path.join(user_home, "projects"),
        os.path.join(user_home, "repos"),
        os.path.join(user_home, "work"),
        user_home,
        os.getcwd(),
    ]
    
    # Filter out non-existent directories and ensure uniqueness
    search_dirs = list(dict.fromkeys(d for d in search_dirs if os.path.exists(d)))
    
    # Search for exact match first
    for directory in search_dirs:
        try:
            for entry in os.listdir(directory):
                full_path = os.path.join(directory, entry)
                if os.path.isdir(full_path) and entry.lower() == dirname.lower():
                    return full_path
        except Exception:
            continue
            
    # Search recursively with limited depth if not found
    for directory in search_dirs:
        try:
            for root, dirs, files in os.walk(directory):
                # Limit depth to avoid performance issues
                try:
                    rel = os.path.relpath(root, directory)
                    depth = 0 if rel == '.' else len(rel.replace('\\', '/').split('/'))
                except ValueError:
                    depth = 0
                if depth >= 4:
                    del dirs[:] # Don't go deeper
                    
                for d in dirs:
                    if d.lower() == dirname.lower() or dirname.lower() in d.lower():
                        return os.path.join(root, d)
        except Exception:
            continue
            
    return None


def find_file_broadly_tool(args: Dict[str, Any]) -> str:
    """Attempt to find a file by name in common user directories."""
    filename = args.get('filename', '')
    if not filename:
        return "Error: No filename provided."
    found_path = find_file_broadly(filename)
    if found_path:
        return f"File found at: {found_path}"
    else:
        return f"Could not find file '{filename}' in common directories."


def find_directory_broadly_tool(args: Dict[str, Any]) -> str:
    """Attempt to find a directory by name in common user directories."""
    dirname = args.get('dirname', '')
    if not dirname:
        return "Error: No directory name provided."
    found_path = find_directory_broadly(dirname)
    if found_path:
        return f"Directory found at: {found_path}"
    else:
        return f"Could not find directory '{dirname}' in common directories."


def search_file_tool(args: Dict[str, Any]) -> str:
    """Search for a file by name across common directories and return its absolute path."""
    filename = args.get('filename', '')
    if not filename:
        return "Error: No filename provided."
        
    found_path = find_file_broadly(filename)
    if found_path:
        return f"File found at: {found_path}"
    else:
        return f"Could not find file '{filename}' in common directories."


def read_file_tool(args: Dict[str, Any]) -> str:
    """Read the contents of a given file path (absolute or relative) with enhanced error handling for binary files and size limits."""
    path = args.get('path', '')
    offset = args.get('offset', 0)
    limit = args.get('limit', 10000)  # Reduced default limit to 10k characters to prevent rate limits and save tokens

    if not path:
        return "No path provided."

    # Try to find the file if it doesn't exist at the given path
    actual_path = path
    if not os.path.exists(path):
        found_path = find_file_broadly(path)
        if found_path:
            actual_path = found_path
        else:
            return f"Error: File '{path}' does not exist and could not be found in common directories."

    if os.path.isdir(actual_path):
        return f"'{actual_path}' is a directory, not a file."

    try:
        # Read initial bytes to determine file type
        with open(actual_path, 'rb') as f:
            initial_bytes = f.read(1024)
            # Character set to check for non-text files
            text_characters = bytearray({7, 8, 9, 10, 12, 13, 27} | set(range(0x20, 0x100)))
            is_binary = bool(initial_bytes.translate(None, text_characters))

            if is_binary:
                # If it's binary, check if it's an image we can recognize
                ext = os.path.splitext(actual_path)[1].lower()
                if ext in ['.jpg', '.jpeg', '.png']:
                    return f"The file '{actual_path}' appears to be an image. Use 'recognize_image' tool to analyze it."
                return f"The file '{actual_path}' appears to be binary and cannot be read as text."

        # If the file is text, proceed to read it with limit and offset
        file_size = os.path.getsize(actual_path)
        with open(actual_path, 'r', encoding='utf-8', errors='replace') as f:
            if offset > 0:
                f.seek(offset)
            content = f.read(limit)
        
        result = f"Content of {actual_path}:\n\n" + content
        if file_size > (offset + limit):
            result += f"\n\n[... Output truncated. File size: {file_size} bytes. Read {limit} characters from offset {offset}. Use 'limit' and 'offset' to read more ...]"
        
        return result
    except IOError as e:
        return f"File access issue: {e}"
    except UnicodeDecodeError:
        return "The file contains non-UTF-8 characters and appears to be binary."
    except Exception as e:
        return f"An unexpected problem occurred: {e}"


def list_files_tool(args: Dict[str, Any]) -> str:
    """List files and directories at a given path (absolute or relative)."""
    try:
        path = args.get('path', '.')
        if not path:
            path = '.'

        actual_path = path
        if not os.path.exists(path):
            found_path = find_directory_broadly(path)
            if found_path:
                actual_path = found_path
            else:
                return f"Error: Path '{path}' does not exist and could not be found in common directories."
        
        if not os.path.isdir(actual_path):
            return f"Error: Path '{actual_path}' is not a directory."

        entries = os.listdir(actual_path)
        files = []
        for entry in entries:
            full_path = os.path.join(actual_path, entry)
            if os.path.isdir(full_path):
                files.append(f"{entry}/")
            else:
                files.append(entry)

        return json.dumps(sorted(files))
    except Exception as e:
        return f"Error: {str(e)}"
    
def playwright_mcp_tool(args):
    import asyncio
    import sys
    from mcp import ClientSession
    from mcp.client.stdio import stdio_client, StdioServerParameters

    async def _call():
        server_params = StdioServerParameters(
            command=sys.executable,
            args=["-B", os.path.join(os.path.dirname(os.path.abspath(__file__)), "mcp_playwright_server.py")],
            cwd=os.path.dirname(os.path.abspath(__file__))
        )
        async with stdio_client(server_params) as (read, write):
            async with ClientSession(read, write) as session:
                await session.initialize()
                result = await session.call_tool("navigate", args)
                return result.content[0].text if result.content else ""

    try:
        return asyncio.run(_call())
    except Exception as e:
        return f"Playwright MCP error: {e}"

def change_working_directory_tool(args: Dict[str, Any]) -> str:
    """Change the current working directory of the agent to a different project or folder."""
    path = args.get('path', '')
    if not path:
        return "Error: No path provided."
    
    actual_path = path
    if not os.path.exists(path):
        found_path = find_directory_broadly(path)
        if found_path:
            actual_path = found_path
        else:
            return f"Error: Directory '{path}' not found and could not be located in common directories."
    
    if not os.path.isdir(actual_path):
        return f"Error: '{actual_path}' is not a directory."
    
    try:
        os.chdir(actual_path)
        return f"Successfully changed working directory to: {os.getcwd()}"
    except Exception as e:
        return f"Error changing directory: {str(e)}"


def create_new_file(file_path: str, content: str) -> str:
    """Create a new file with the given content at any path (absolute or relative)."""
    try:
        # Create directory if it doesn't exist
        directory = os.path.dirname(file_path)
        if directory and directory != '.':
            os.makedirs(directory, exist_ok=True)

        # Write the file
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(content)

        return f"Successfully created file {file_path}"
    except Exception as e:
        return f"Error: {str(e)}"


def delete_file_tool(args: Dict[str, Any]) -> str:
    """Delete a file or directory at any path (absolute or relative)."""
    try:
        path = args.get('path', '')
        if not path:
            return "Error: no path provided"
            
        actual_path = path
        if not os.path.exists(path):
            # Try to find a directory first if the user mentions "folder" or if it ends in /
            found_path = None
            if path.endswith('/') or path.endswith('\\'):
                found_path = find_directory_broadly(path.rstrip('/\\'))
            
            if not found_path:
                found_path = find_directory_broadly(path)
            if not found_path:
                found_path = find_file_broadly(path)
            
            if found_path:
                actual_path = found_path
            else:
                return f"Error: '{path}' does not exist and could not be found in common directories."

        if os.path.isdir(actual_path):
            shutil.rmtree(actual_path)
            return f"Successfully deleted directory {actual_path}"
        else:
            os.remove(actual_path)
            return f"Successfully deleted file {actual_path}"
    except Exception as e:
        return f"Error: {str(e)}"


def create_and_edit_file_tool(args: Dict[str, Any]) -> str:
    """Create or edit any file at any path (absolute or relative) by replacing old_str with new_str, or writing new_str if the file does not exist."""
    try:
        path = args.get('path', '')
        old_str = args.get('old_str', '').replace('\\n', '\n')
        new_str = args.get('new_str', '').replace('\\n', '\n')

        if not path:
            return "Error: No path provided."
        
        actual_path = path
        if not os.path.exists(path):
            # If path doesn't exist, check if the parent directory exists or can be found
            parent_dir = os.path.dirname(path)
            filename = os.path.basename(path)
            
            if parent_dir and not os.path.exists(parent_dir):
                found_parent = find_directory_broadly(parent_dir)
                if found_parent:
                    actual_path = os.path.join(found_parent, filename)
        
        if old_str == new_str:
            return "Error: old_str and new_str are the same; no changes made."

        # Try to read the existing file
        if not os.path.exists(actual_path):
            if old_str == "":
                return create_new_file(actual_path, new_str)
            else:
                return f"Error: File '{actual_path}' not found and no content provided to create it."

        with open(actual_path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()

        # Replace old_str with new_str
        if old_str not in content and old_str != "":
            return f"Error: old_str not found in '{actual_path}'"

        new_content = content.replace(old_str, new_str) if old_str != "" else new_str

        # Write the modified content back to file
        with open(actual_path, 'w', encoding='utf-8') as f:
            f.write(new_content)
        
        return f"Successfully updated '{actual_path}'"

    except Exception as e:
        return f"Error editing file: {str(e)}"


def rename_file_tool(args: Dict[str, Any]) -> str:
    """Rename a file at any path (absolute or relative). Args: old_path (str), new_path (str)."""
    try:
        old_path = args.get('old_path', '')
        new_path = args.get('new_path', '')
        if not old_path or not new_path:
            return "Error: old_path and new_path are required."
        if not os.path.exists(old_path):
            return f"Error: file {old_path} does not exist."
        if os.path.exists(new_path):
            return f"Error: file {new_path} already exists."
        os.rename(old_path, new_path)
        return f"Successfully renamed {old_path} to {new_path}"
    except Exception as e:
        return f"Error: {str(e)}"


def run_code_tool(args: Dict[str, Any]) -> str:
    """Execute code locally."""
    import subprocess
    import os
    try:
        command = args.get('command', '')
        if not command:
            return "Error: no command provided"
        
        print(f"\033[92mExecuting Command:\033[0m {command}")

        result = subprocess.run(
            command,
            shell=True,
            capture_output=True,
            text=True
        )
        
        stdout = result.stdout
        stderr = result.stderr
        
        limit = 10000
        if len(stdout) > limit:
            stdout = stdout[:limit] + f"\n\n[... STDOUT truncated. Total size: {len(result.stdout)} characters ...]"
        if len(stderr) > limit:
            stderr = stderr[:limit] + f"\n\n[... STDERR truncated. Total size: {len(result.stderr)} characters ...]"
            
        output = f"STDOUT:\n{stdout}\nSTDERR:\n{stderr}"
        return output
    except subprocess.TimeoutExpired:
        return "Error: Command timed out."
    except Exception as e:
        return f"Error: {str(e)}"


def check_syntax_tool(args: Dict[str, Any]) -> str:
    """Check the syntax of a code file."""
    path = args.get('path', '')
    if not path:
        return "Error: no path provided"
    
    ext = os.path.splitext(path)[1].lower()
    if ext == '.py':
        command = f"python3 -m py_compile {path}"
    elif ext == '.java':
        command = f"javac {path}"
    elif ext in ['.c']:
        command = f"gcc -fsyntax-only {path}"
    elif ext in ['.cpp', '.cc', '.cxx']:
        command = f"g++ -fsyntax-only {path}"
    elif ext == '.rs':
        if os.path.exists('Cargo.toml'):
            command = "cargo check"
        else:
            command = f"rustc --edition 2021 -o /dev/null --emit=dep-info {path}"
    elif ext in ['.js', '.jsx']:
        command = f"node --check {path}"
    elif ext in ['.ts', '.tsx']:
        command = f"tsc --noEmit {path}"
    elif ext == '.go':
        command = f"go build -o /dev/null {path}"
    elif ext == '.sql':
        command = f"sqlfluff lint {path} --dialect ansi"
    else:
        return f"Error: syntax check not supported for extension {ext}"
    
    return run_code_tool({'command': command})


def run_tests_tool(args: Dict[str, Any]) -> str:
    """Run tests for the project."""
    command = args.get('command', '')
    if not command:
        # Try to infer command
        if os.path.exists('pytest.ini') or os.path.exists('tests/'):
            command = "pytest"
        elif os.path.exists('Cargo.toml'):
            command = "cargo test"
        elif os.path.exists('package.json'):
            command = "npm test"
        elif os.path.exists('go.mod'):
            command = "go test ./..."
        elif os.path.exists('pom.xml'):
            command = "mvn test"
        elif os.path.exists('build.gradle'):
            command = "gradle test"
        elif os.path.exists('Makefile'):
            command = "make test"
        else:
            return "Error: no test command provided and could not infer one"
    
    return run_code_tool({'command': command})

def _github_mcp_call(tool_name, args):
    import asyncio
    import sys
    from mcp import ClientSession
    from mcp.client.stdio import stdio_client, StdioServerParameters

    async def _call():
        server_params = StdioServerParameters(
            command=sys.executable,
            args=["-B", os.path.join(os.path.dirname(os.path.abspath(__file__)), "mcp_github_server.py")],
            cwd=os.path.dirname(os.path.abspath(__file__))
        )
        async with stdio_client(server_params) as (read, write):
            async with ClientSession(read, write) as session:
                await session.initialize()
                result = await session.call_tool(tool_name, args)
                return result.content[0].text if result.content else ""

    try:
        return asyncio.run(_call())
    except Exception as e:
        return f"GitHub MCP error: {e}"


def github_create_pr_tool(args):
    return _github_mcp_call("create_pull_request", args)


def github_create_branch_tool(args):
    return _github_mcp_call("create_branch", args)


def github_commit_file_tool(args):
    return _github_mcp_call("commit_file", args)


def github_commit_local_file_tool(args):
    return _github_mcp_call("commit_local_file", args)


def create_github_issue_tool(args: Dict[str, Any]) -> str:
    """Create a GitHub issue on the configured repository."""
    try:
        title = args.get('title', '')
        body = args.get('body', '')
        labels = args.get('labels', [])
        assignees = args.get('assignees', [])

        if not title:
            return "Error: Issue title is required."

        # Get GitHub credentials from settings
        github_token = getattr(settings, 'GITHUB_TOKEN', None)
        github_repo = getattr(settings, 'GITHUB_REPO', None)

        if not github_token or not github_repo:
            return "Error: GITHUB_TOKEN or GITHUB_REPO not configured in settings."

        # GitHub API endpoint
        url = f"https://api.github.com/repos/{github_repo}/issues"
        headers = {
            "Authorization": f"token {github_token}",
            "Accept": "application/vnd.github.v3+json"
        }

        payload = {
            "title": title,
            "body": body
        }

        if labels:
            payload["labels"] = labels if isinstance(labels, list) else [labels]
        if assignees:
            payload["assignees"] = assignees if isinstance(assignees, list) else [assignees]

        response = requests.post(url, headers=headers, json=payload)

        if response.status_code == 201:
            issue_data = response.json()
            issue_number = issue_data.get('number')
            issue_url = issue_data.get('html_url')
            return f"Successfully created issue #{issue_number}: {issue_url}"
        else:
            return f"Error creating issue: {response.status_code} - {response.text}"

    except Exception as e:
        return f"Error: {str(e)}"


def lint_code_tool(args: Dict[str, Any]) -> str:
    """Run static analysis (linting) on a code file."""
    path = args.get('path', '')
    if not path:
        return "Error: no path provided"
    
    ext = os.path.splitext(path)[1].lower()
    if ext == '.py':
        command = f"pylint {path}"
    elif ext == '.java':
        # Simple checkstyle-like check with javac warnings
        command = f"javac -Xlint:all {path}"
    elif ext in ['.c', '.cpp', '.cc', '.cxx']:
        command = f"cppcheck {path}"
    elif ext == '.rs':
        command = "cargo clippy" if os.path.exists('Cargo.toml') else f"rustc -W help"
    elif ext in ['.js', '.jsx', '.ts', '.tsx']:
        command = f"eslint {path}"
    elif ext == '.go':
        command = f"go vet {path}"
    elif ext == '.sql':
        command = f"sqlfluff lint {path} --dialect ansi"
    else:
        return f"Error: linting not supported for extension {ext}"
    
    return run_code_tool({'command': command})


def create_gmail_draft(recipient: str, subject: str, body: str, attachments: List[str]) -> str:
    """Create a draft in Gmail via IMAP."""
    user = getattr(settings, 'GMAIL_SENDER_ADDRESS', None)
    password = getattr(settings, 'GMAIL_APP_PASSWORD', None)

    if not user or not password:
        return "Error: GMAIL_SENDER_ADDRESS or GMAIL_APP_PASSWORD not configured in settings."

    try:
        # Create message with robust structure
        msg = MIMEMultipart('mixed')
        msg['To'] = recipient
        msg['From'] = user
        msg['Subject'] = subject
        msg['Date'] = formatdate(localtime=True)
        msg['Message-ID'] = make_msgid()
        
        # Create the body container (multipart/alternative)
        body_part = MIMEMultipart('alternative')
        
        # Add plain text version
        body_part.attach(MIMEText(body, 'plain'))
        
        # Add HTML version (simple conversion)
        html_body = body.replace('\n', '<br>')
        body_part.attach(MIMEText(f"<html><body>{html_body}</body></html>", 'html'))
        
        # Attach body container to the main message
        msg.attach(body_part)

        for path in attachments:
            if os.path.exists(path):
                filename = os.path.basename(path)
                part = MIMEBase('application', "octet-stream")
                try:
                    with open(path, 'rb') as file:
                        part.set_payload(file.read())
                    encoders.encode_base64(part)
                    part.add_header('Content-Disposition', f'attachment; filename="{filename}"')
                    msg.attach(part)
                except Exception as file_err:
                    return f"Error: Reading attachment {path} failed: {str(file_err)}"
            else:
                return f"Error: Attachment file not found at {path}"

        # Connect to Gmail IMAP
        try:
            imap = imaplib.IMAP4_SSL("imap.gmail.com")
            imap.login(user, password)
        except imaplib.IMAP4.error as auth_err:
            return f"Error: Authentication failed. Please ensure: 1. Your GMAIL_PASSWORD in .env is a 16-character App Password. 2. IMAP is enabled. Original error: {str(auth_err)}"
        except Exception as conn_err:
            return f"Error: Failed to connect to Gmail IMAP: {str(conn_err)}"

        # Try to find the Drafts folder
        status, folders = imap.list()
        draft_folder = None
        if status == 'OK':
            for folder in folders:
                folder_str = folder.decode('utf-8')
                # Gmail drafts folder has the \Drafts attribute
                if '\\Drafts' in folder_str:
                    # The folder name is usually the last quoted string
                    matches = re.findall(r'"([^"]+)"', folder_str)
                    if matches:
                        draft_folder = matches[-1]
                    break
        
        if not draft_folder:
            # Fallback for some configurations
            draft_folder = "[Gmail]/Drafts"

        # Append to drafts
        try:
            # Ensure draft_folder is quoted if it contains spaces or special characters
            quoted_folder = draft_folder
            if not draft_folder.startswith('"') and any(c in draft_folder for c in ' []/'):
                quoted_folder = f'"{draft_folder}"'

            # Use CRLF for line endings as required by IMAP
            msg_bytes = msg.as_bytes(policy=SMTP)
            res, detail = imap.append(quoted_folder, r'(\Draft)', imaplib.Time2Internaldate(time.time()), msg_bytes)
            if res != 'OK':
                return f"Error: IMAP APPEND failed: {res} - {str(detail)}"
        except Exception as append_err:
            return f"Error: Exception during IMAP APPEND to '{draft_folder}': {str(append_err)}"
            
        imap.logout()
        return "OK"
    except Exception as e:
        return f"Error: {str(e)}"


def find_chrome_profile_for_email(email: str) -> str:
    """Attempt to find the Chrome profile name associated with a specific email."""
    if not email:
        return "Default"
    
    system = platform.system()
    if system == "Windows":
        base_path = os.path.expandvars(r"%LocalAppData%\Google\Chrome\User Data")
    elif system == "Darwin":
        base_path = os.path.expanduser("~/Library/Application Support/Google/Chrome")
    else:
        base_path = os.path.expanduser("~/.config/google-chrome")

    if not os.path.exists(base_path):
        return "Default"

    try:
        profiles = [d for d in os.listdir(base_path) if os.path.isdir(os.path.join(base_path, d)) and (d == "Default" or d.startswith("Profile"))]
        
        for profile in profiles:
            pref_path = os.path.join(base_path, profile, "Preferences")
            if os.path.exists(pref_path):
                try:
                    with open(pref_path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                        if email.lower() in content.lower():
                            return profile
                except Exception:
                    continue
    except Exception:
        pass
        
    return "Default"


def open_url_in_chrome_profile(url: str, email: str = None) -> bool:
    """Attempt to open a URL in a specific Chrome profile."""
    # Priority: 1. Explicitly configured directory, 2. Auto-discovered directory, 3. Default
    profile = getattr(settings, 'CHROME_PROFILE_DIRECTORY', None)
    if not profile or profile == "Default":
        profile = find_chrome_profile_for_email(email)
    
    system = platform.system()
    
    try:
        if system == "Windows":
            # Common paths for Chrome on Windows
            chrome_paths = [
                os.path.expandvars(r"%ProgramFiles%\Google\Chrome\Application\chrome.exe"),
                os.path.expandvars(r"%ProgramFiles(x86)%\Google\Chrome\Application\chrome.exe"),
                os.path.expandvars(r"%LocalAppData%\Google\Chrome\Application\chrome.exe")
            ]
            for path in chrome_paths:
                if os.path.exists(path):
                    subprocess.Popen([path, f"--profile-directory={profile}", url])
                    return True
        elif system == "Darwin":  # macOS
            subprocess.Popen(["open", "-a", "Google Chrome", "--args", f"--profile-directory={profile}", url])
            return True
        elif system == "Linux":
            subprocess.Popen(["google-chrome", f"--profile-directory={profile}", url])
            return True
    except Exception:
        pass
    
    # Fallback to default browser
    return webbrowser.open_new_tab(url)


def open_gmail_and_compose_tool(args: Dict[str, Any]) -> str:
    user = getattr(settings, 'GMAIL_SENDER_ADDRESS', '')
    recipient = args.get('recipient', '').strip()
    subject = args.get('subject', '')
    body = args.get('body', '')
    attachments = args.get('attachments', [])
    if not recipient:
        return "Error: recipient is required"

    if isinstance(attachments, str):
        attachments = [attachments]

    # Use user-specific URL if available
    base_url = f"https://mail.google.com/mail/u/{user}/" if user else "https://mail.google.com/mail/"

    # If there are attachments, we create a draft via IMAP because web URL doesn't support them
    if attachments:
        draft_result = create_gmail_draft(recipient, subject, body, attachments)
        if draft_result == "OK":
            drafts_url = f"{base_url}#drafts"
            try:
                opened = open_url_in_chrome_profile(drafts_url, user)
                if opened:
                    return f"A draft with the attachments has been created in your Gmail Drafts ({user}). I've opened your Drafts folder in the browser (Profile: {find_chrome_profile_for_email(user)}). Please review and send it."
                return f"A draft with the attachments has been created in your Gmail Drafts ({user}). Please open {drafts_url} to review and send it."
            except:
                return f"A draft with the attachments has been created in your Gmail Drafts ({user}). Please open {drafts_url} to review and send it."
        else:
            # If draft fails, we still want to open the compose window as a fallback
            fallback_msg = f"Failed to create draft with attachments: {draft_result}."
            
            query = [f"to={quote(recipient)}"]
            if subject:
                query.append(f"su={quote(subject)}")
            if body:
                query.append(f"body={quote(body)}")
            compose_url = f"{base_url}?view=cm&fs=1&tf=1&" + "&".join(query)
            
            try:
                opened = open_url_in_chrome_profile(compose_url, user)
                if opened:
                    return f"{fallback_msg} I've opened the standard compose window for you instead (Profile: {find_chrome_profile_for_email(user)}). Please attach the files manually and send."
                return f"{fallback_msg} Please open this link to compose manually: {compose_url}"
            except:
                return f"{fallback_msg} Please open this link to compose manually: {compose_url}"

    query = [f"to={quote(recipient)}"]
    if subject:
        query.append(f"su={quote(subject)}")
    if body:
        query.append(f"body={quote(body)}")
    compose_url = f"{base_url}?view=cm&fs=1&tf=1&" + "&".join(query)

    try:
        opened = open_url_in_chrome_profile(compose_url, user)
        if opened:
            return f"Gmail compose window for {user} opened in your browser (Profile: {find_chrome_profile_for_email(user)}). Log in if prompted, review the message, and press Send."
        return f"Open this link manually to compose the email: {compose_url}"
    except Exception as e:
        return f"Error launching browser automatically ({str(e)}). Open this link manually: {compose_url}"


def recognize_image_tool(args: Dict[str, Any]) -> str:
    """Use GPT-4o vision to recognize contents of an image (jpg, png)."""
    path = args.get('path', '')
    prompt = args.get('prompt', 'What is in this image? Provide a detailed description.')

    if not path or not os.path.exists(path):
        return f"Error: File '{path}' not found."

    try:
        # Read and encode image to base64
        with open(path, "rb") as image_file:
            base64_image = base64.b64encode(image_file.read()).decode('utf-8')

        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        },
                    ],
                }
            ],
            max_tokens=500,
        )

        return response.choices[0].message.content
    except Exception as e:
        return f"Error during image recognition: {str(e)}"


def recognize_video_tool(args: Dict[str, Any]) -> str:
    """Use GPT-4o vision to recognize contents of a video (mp4) by extracting frames."""
    path = args.get('path', '')
    prompt = args.get('prompt', 'These are frames from a video. What is happening? Provide a summary and details.')
    max_frames = args.get('max_frames', 10)

    if not path or not os.path.exists(path):
        return f"Error: File '{path}' not found."

    try:
        video = cv2.VideoCapture(path)
        total_frames = int(video.get(cv2.CAP_PROP_FRAME_COUNT))
        fps = video.get(cv2.CAP_PROP_FPS)
        
        if total_frames <= 0:
            return "Error: Could not read frames from video."

        # Sample frames evenly
        interval = max(1, total_frames // max_frames)
        base64_frames = []
        
        count = 0
        while video.isOpened():
            success, frame = video.read()
            if not success or len(base64_frames) >= max_frames:
                break
            
            if count % interval == 0:
                _, buffer = cv2.imencode(".jpg", frame)
                base64_frames.append(base64.b64encode(buffer).decode("utf-8"))
            
            count += 1
        
        video.release()

        if not base64_frames:
            return "Error: No frames extracted from video."

        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        
        content = [{"type": "text", "text": prompt}]
        for base64_frame in base64_frames:
            content.append({
                "type": "image_url",
                "image_url": {
                    "url": f"data:image/jpeg;base64,{base64_frame}"
                }
            })

        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": content,
                }
            ],
            max_tokens=1000,
        )

        return response.choices[0].message.content
    except Exception as e:
        return f"Error during video recognition: {str(e)}"


def recognize_audio_tool(args: Dict[str, Any]) -> str:
    """Use GPT-4o to analyze and summarize an audio file, including speech, music, and sounds."""
    path = args.get('path', '')
    prompt = args.get('prompt', 'Analyze this audio file. Describe everything you hear: transcribe any speech, identify any music (genre, instruments, mood), and note any ambient or other sounds.')

    SUPPORTED_FORMATS = ('.wav', '.mp3', '.ogg', '.flac', '.webm', '.m4a', '.mp4', '.aac', '.wma', '.opus')

    if not path or not os.path.exists(path):
        return f"Error: File '{path}' not found."

    ext = os.path.splitext(path)[1].lower()
    if ext not in SUPPORTED_FORMATS:
        return f"Error: Unsupported audio format '{ext}'. Supported formats: {', '.join(SUPPORTED_FORMATS)}"

    converted_path = None
    try:
        # GPT-4o audio input only accepts wav and mp3; convert other formats
        if ext in ('.wav', '.mp3'):
            audio_path = path
            audio_format = ext[1:]  # strip the leading dot
        else:
            from pydub import AudioSegment
            audio_segment = AudioSegment.from_file(path)
            converted_path = os.path.splitext(path)[0] + '_temp_converted.wav'
            audio_segment.export(converted_path, format='wav')
            audio_path = converted_path
            audio_format = 'wav'

        with open(audio_path, "rb") as f:
            base64_audio = base64.b64encode(f.read()).decode('utf-8')

        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        response = client.chat.completions.create(
            model="gpt-4o-audio-preview",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "input_audio",
                            "input_audio": {
                                "data": base64_audio,
                                "format": audio_format
                            }
                        }
                    ],
                }
            ],
            max_tokens=1000,
        )

        return response.choices[0].message.content
    except Exception as e:
        return f"Error during audio recognition: {str(e)}"
    finally:
        if converted_path and os.path.exists(converted_path):
            os.remove(converted_path)

SEARCH_FILE_DEFINITION = ToolDefinition(
    name="search_file",
    description="Search for a file by name across common directories (Desktop, Documents, Downloads, Pictures, etc.) and return its absolute path.",
    parameters={
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "The name of the file to search for (e.g., 'my meetings.png')."
            }
        },
        "required": ["filename"]
    },
    function=search_file_tool
)


# Define the read_file tool
READ_FILE_DEFINITION = ToolDefinition(
    name="read_file",
    description="Read the contents of a given file path (absolute or relative). Use this when you want to see what's inside a file. Do not use this with directory names.",
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The absolute or relative path of a file."
            },
            "offset": {
                "type": "integer",
                "description": "The character offset to start reading from."
            },
            "limit": {
                "type": "integer",
                "description": "The maximum number of characters to read."
            }
        },
        "required": ["path"]
    },
    function=read_file_tool
)

# Define the list_files tool
LIST_FILES_DEFINITION = ToolDefinition(
    name="list_files",
    description="List files and directories at a given path (absolute or relative). If no path is provided, lists files in the current directory.",
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "Optional absolute or relative path to list files from. Defaults to current directory if not provided."
            }
        },
        "required": []
    },
    function=list_files_tool
)


DELETE_FILE_DEFINITION = ToolDefinition(
    name="delete_file",
    description="Delete a file or directory at any path (absolute or relative).",
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The absolute or relative path of the file or directory to delete."
            }
        },
        "required": ["path"]
    },
    function=delete_file_tool,
    requires_approval=True
)


# Define the edit_file tool
CREATE_AND_EDIT_FILE_DEFINITION = ToolDefinition(
    name="create_and_edit_file",
    description="""Create or edit any file at any path (absolute or relative) by replacing old_str with new_str, or writing new_str if the file does not exist. Supports all file types, including code files such as .cpp, .py, .js, etc. IMPORTANT: Always provide the full code with proper indentation and newlines. Avoid writing code in a single line.""",
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The absolute or relative path to the file (any file type allowed)"
            },
            "old_str": {
                "type": "string",
                "description": "Text to search for - must match exactly and must only have one match exactly. If you want to replace the entire file content, you can first read it and then provide it here, or use old_str='' to overwrite/create."
            },
            "new_str": {
                "type": "string",
                "description": "Text to replace old_str with, or to write if creating a new file. Ensure you use actual newline characters (\\n) and proper indentation for code files."
            }
        },
        "required": ["path", "old_str", "new_str"]
    },
    function=create_and_edit_file_tool,
    requires_approval=True
)

GITHUB_CREATE_BRANCH_DEFINITION = ToolDefinition(
    name="github_create_branch",
    description="Step 1 of 3: Create a new branch on the configured GitHub repository. Call this FIRST before committing files or creating a PR. Do NOT use run_code with git for GitHub operations — use these MCP tools instead.",
    parameters={
        "type": "object",
        "properties": {
            "name": {"type": "string", "description": "Name of the new branch"},
            "source": {"type": "string", "description": "Branch to create from (default: main)"}
        },
        "required": ["name"]
    },
    function=github_create_branch_tool,
    requires_approval=True
)

GITHUB_COMMIT_FILE_DEFINITION = ToolDefinition(
    name="github_commit_file",
    description="Step 2 of 3: Commit a file with given content to a branch on the configured GitHub repository. Use when file content is provided in the prompt. Call AFTER github_create_branch, BEFORE github_create_pr.",
    parameters={
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "File path in the repo"},
            "content": {"type": "string", "description": "Full file content to commit"},
            "message": {"type": "string", "description": "Commit message"},
            "branch": {"type": "string", "description": "Target branch"}
        },
        "required": ["path", "content", "message", "branch"]
    },
    function=github_commit_file_tool,
    requires_approval=True
)

GITHUB_COMMIT_LOCAL_FILE_DEFINITION = ToolDefinition(
    name="github_commit_local_file",
    description="Step 2 of 3: Read a local file and commit it to a branch on the configured GitHub repository. Use when committing an existing local file. Call AFTER github_create_branch, BEFORE github_create_pr.",
    parameters={
        "type": "object",
        "properties": {
            "local_path": {"type": "string", "description": "Local file path to read"},
            "message": {"type": "string", "description": "Commit message"},
            "branch": {"type": "string", "description": "Target branch"},
            "remote_path": {"type": "string", "description": "Path in repo (defaults to filename)"}
        },
        "required": ["local_path", "message", "branch"]
    },
    function=github_commit_local_file_tool,
    requires_approval=True
)

GITHUB_MCP_DEFINITION = ToolDefinition(
    name="github_create_pr",
    description="Step 3 of 3: Create a GitHub pull request on the configured repository. ONLY call this AFTER a branch has been created (github_create_branch) and at least one file committed to it (github_commit_file or github_commit_local_file). Do NOT use run_code with git for any GitHub operations.",
    parameters={
        "type": "object",
        "properties": {
            "title": {"type": "string", "description": "Title for the pull request"},
            "head": {"type": "string", "description": "The source branch with changes"},
            "base": {"type": "string", "description": "The target branch (default: main)"},
            "body": {"type": "string", "description": "Pull request description"}
        },
        "required": ["title", "head"]
    },
    function=github_create_pr_tool,
    requires_approval=True
)


CREATE_GITHUB_ISSUE_DEFINITION = ToolDefinition(
    name="create_github_issue",
    description="Create a GitHub issue on the configured repository. Use this to track bugs, feature requests, tasks, or automatically create tickets when code execution fails or tests fail. Issues can be labeled and assigned.",
    parameters={
        "type": "object",
        "properties": {
            "title": {
                "type": "string",
                "description": "The title of the issue (required). Keep it concise and descriptive."
            },
            "body": {
                "type": "string",
                "description": "The body/description of the issue. Supports markdown formatting. Include error details, stack traces, reproduction steps, etc."
            },
            "labels": {
                "type": "array",
                "items": {"type": "string"},
                "description": "Optional array of label names to apply (e.g., ['bug', 'high-priority', 'test-failure'])."
            },
            "assignees": {
                "type": "array",
                "items": {"type": "string"},
                "description": "Optional array of GitHub usernames to assign the issue to."
            }
        },
        "required": ["title"]
    },
    function=create_github_issue_tool,
    requires_approval=True
)


PLAYWRIGHT_MCP_DEFINITION = ToolDefinition(
    name="playwright_navigate",
    description="Navigate websites using Playwright MCP server",
    parameters={
        "type": "object",
        "properties": {
            "url": {"type": "string"},
            "screenshot": {"type": "string"}
        },
        "required": ["url"]
    },
    function=playwright_mcp_tool,
    requires_approval=True
)


# ─── Travel Search Tools ─────────────────────────────────────────────

def search_flights_tool(args: Dict[str, Any]) -> str:
    """Search for flights using the Duffel REST API."""

    origin = args.get('origin', '').strip().upper()
    destination = args.get('destination', '').strip().upper()
    departure_date = args.get('departure_date', '').strip()
    return_date = args.get('return_date', '').strip() if args.get('return_date') else ''
    adults = int(args.get('adults', 1))
    cabin_class = args.get('cabin_class', 'economy').strip().lower()

    if not origin or not destination or not departure_date:
        return "Error: origin, destination, and departure_date are required."

    try:
        # Build request payload
        slices = [{"origin": origin, "destination": destination, "departure_date": departure_date}]
        if return_date:
            slices.append({"origin": destination, "destination": origin, "departure_date": return_date})

        payload = {
            "data": {
                "slices": slices,
                "passengers": [{"type": "adult"} for _ in range(adults)],
                "cabin_class": cabin_class,
                "return_offers": True,
            }
        }

        status, resp_json = _duffel_post("/air/offer_requests", payload, timeout=45)

        if status not in (200, 201):
            error_msg = resp_json.get("errors", [{}])[0].get("message", str(resp_json)) if resp_json.get("errors") else str(resp_json)
            return f"Flight search failed (HTTP {status}): {error_msg}"

        data = resp_json.get("data", {})
        offers = data.get("offers", [])

        if not offers:
            return (
                f"No flights found from {origin} to {destination} on {departure_date}. "
                f"No offers are available for these dates/route."
            )

        # Cache passenger IDs for later booking
        passenger_ids = [p["id"] for p in data.get("passengers", [])]

        # Sort offers by price
        sorted_offers = sorted(offers, key=lambda o: float(o.get("total_amount", "9999")))

        lines = [
            f"Found {len(sorted_offers)} flight offer(s) from {origin} to "
            f"{destination} on {departure_date}:\n"
        ]

        for idx, offer in enumerate(sorted_offers[:15], 1):
            offer_id = offer.get("id", "")
            airline = offer.get("owner", {}).get("name", "Unknown Airline")
            total = offer.get("total_amount", "?")
            currency = offer.get("total_currency", "USD")

            # Cache this offer for booking
            _duffel_offer_cache[offer_id] = {
                "passenger_ids": passenger_ids,
                "total_amount": total,
                "total_currency": currency,
                "expires_at": offer.get("expires_at", ""),
            }

            # Extract first slice details for display
            slices_data = offer.get("slices", [])
            first_slice = slices_data[0] if slices_data else {}

            dep_apt = first_slice.get("origin", {}).get("iata_code", origin)
            arr_apt = first_slice.get("destination", {}).get("iata_code", destination)
            duration = _parse_iso_duration(first_slice.get("duration", ""))
            segments = first_slice.get("segments", [])
            num_stops = max(0, len(segments) - 1)
            stops_text = 'Nonstop' if num_stops == 0 else f'{num_stops} stop{"s" if num_stops > 1 else ""}'

            if segments:
                dep_time_raw = segments[0].get("departing_at", "")
                arr_time_raw = segments[-1].get("arriving_at", "")
                dep_time = dep_time_raw[11:16] if len(dep_time_raw) > 16 else ''
                arr_time = arr_time_raw[11:16] if len(arr_time_raw) > 16 else ''
            else:
                dep_time = arr_time = ''

            lines.append(
                f"{idx}. **{airline}** | {dep_apt} {dep_time} \u2192 {arr_apt} {arr_time} | "
                f"{duration} | {stops_text} | **{total} {currency}** "
                f"[Offer: {offer_id}]"
            )

        expires = sorted_offers[0].get("expires_at", "")
        if expires:
            lines.append(f"\n*Offers expire at {expires}. Book promptly to secure these prices.*")

        return '\n'.join(lines)

    except Exception as exc:
        return f"Flight search failed: {exc}"


SEARCH_FLIGHTS_DEFINITION = ToolDefinition(
    name="search_flights",
    description=(
        "Search for flight offers between airports using the Duffel API. Returns a list of available "
        "flights with airlines, times, duration, stops, prices, and offer IDs for booking. "
        "Use IATA airport codes (e.g. LAX, JFK, SYD). Dates must be in YYYY-MM-DD format."
    ),
    parameters={
        "type": "object",
        "properties": {
            "origin": {
                "type": "string",
                "description": "Origin IATA airport code (e.g. 'LAX', 'SYD', 'LHR')."
            },
            "destination": {
                "type": "string",
                "description": "Destination IATA airport code (e.g. 'JFK', 'NRT', 'CDG')."
            },
            "departure_date": {
                "type": "string",
                "description": "Departure date in YYYY-MM-DD format."
            },
            "return_date": {
                "type": "string",
                "description": "Optional return date in YYYY-MM-DD format for round trips."
            },
            "adults": {
                "type": "integer",
                "description": "Number of adult passengers (default 1)."
            },
            "cabin_class": {
                "type": "string",
                "description": "Cabin class: 'economy', 'premium_economy', 'business', or 'first' (default 'economy')."
            }
        },
        "required": ["origin", "destination", "departure_date"]
    },
    function=search_flights_tool,
    requires_approval=False
)


def book_travel_tool(args: Dict[str, Any]) -> str:
    """Book a flight via the Duffel API and persist to database."""
    from chat.models import Booking

    given_name = args.get('given_name', '')
    family_name = args.get('family_name', '')
    passenger_email = args.get('passenger_email', '')
    phone_number = args.get('phone_number', '')
    date_of_birth = args.get('date_of_birth', '')
    gender = args.get('gender', 'm')
    title = args.get('title', 'mr')
    card_last_four = args.get('card_last_four', '')
    card_holder_name = args.get('card_holder_name', '')

    # Duffel identifiers
    offer_id = args.get('offer_id', '')

    # Flight display fields
    airline = args.get('airline', '')
    origin = args.get('origin', '')
    destination = args.get('destination', '')
    departure_date = args.get('departure_date', '')
    departure_time = args.get('departure_time', '')
    arrival_time = args.get('arrival_time', '')
    stops = args.get('stops', 0)
    duration = args.get('duration', '')
    booking_class = args.get('booking_class', 'economy')

    # Validation
    if not given_name or not family_name:
        return "Error: both given_name and family_name are required."
    if not passenger_email:
        return "Error: passenger_email is required."
    if not card_last_four:
        return "Error: card_last_four is required for payment verification."
    if not offer_id:
        return "Error: offer_id is required for flight booking. Use search_flights first to get offer IDs."

    passenger_name = f"{given_name} {family_name}"

    cached = _duffel_offer_cache.get(offer_id)
    if not cached:
        return (
            "Error: Offer not found in cache. The offer may have expired. "
            "Please run search_flights again to get fresh offers."
        )

    try:
        passenger_data = [{
            "id": cached["passenger_ids"][0],
            "given_name": given_name,
            "family_name": family_name,
            "born_on": date_of_birth or "1990-01-01",
            "title": title.lower().rstrip('.'),
            "gender": gender[0].lower() if gender else "m",
            "email": passenger_email,
            "phone_number": phone_number or "+10000000000",
        }]

        # Add additional passengers if there are more cached passenger IDs
        for pid in cached["passenger_ids"][1:]:
            passenger_data.append({
                "id": pid,
                "given_name": given_name,
                "family_name": family_name,
                "born_on": date_of_birth or "1990-01-01",
                "title": title.lower().rstrip('.'),
                "gender": gender[0].lower() if gender else "m",
                "email": passenger_email,
                "phone_number": phone_number or "+10000000000",
            })

        order_payload = {
            "data": {
                "selected_offers": [offer_id],
                "passengers": passenger_data,
                "payments": [{
                    "type": "balance",
                    "currency": cached["total_currency"],
                    "amount": cached["total_amount"],
                }],
                "type": "instant",
            }
        }

        status, resp_json = _duffel_post("/air/orders", order_payload, timeout=45)

        if status not in (200, 201):
            error_msg = resp_json.get("errors", [{}])[0].get("message", str(resp_json)) if resp_json.get("errors") else str(resp_json)
            return f"Flight booking failed (HTTP {status}): {error_msg}"

        order_data = resp_json.get("data", {})
        booking_ref = order_data.get("booking_reference", "") or uuid.uuid4().hex[:10].upper()
        duffel_order_id = order_data.get("id", "")
        total_amount = order_data.get("total_amount", cached["total_amount"])
        total_currency = order_data.get("total_currency", cached["total_currency"])

        # Persist to database
        Booking.objects.create(
            booking_ref=booking_ref,
            booking_type='flight',
            duffel_order_id=duffel_order_id,
            passenger_name=passenger_name,
            passenger_email=passenger_email,
            total_price=Decimal(str(total_amount)),
            currency=total_currency,
            status='confirmed',
            details={
                "duffel_order_id": duffel_order_id,
                "offer_id": offer_id,
                "airline": airline,
                "origin": origin,
                "destination": destination,
                "departure_date": departure_date,
                "departure_time": departure_time,
                "arrival_time": arrival_time,
                "stops": stops,
                "duration": duration,
                "booking_class": booking_class,
                "card_last_four": card_last_four,
            },
        )

        # Remove from cache after successful booking
        _duffel_offer_cache.pop(offer_id, None)

        # Build confirmation
        lines = [
            "## Flight Booking Confirmed!",
            "",
            f"**Booking Reference (PNR):** `{booking_ref}`",
            f"**Duffel Order ID:** `{duffel_order_id}`",
            f"**Status:** Confirmed",
            f"**Passenger:** {passenger_name} ({passenger_email})",
            f"**Payment:** Card ending in ****{card_last_four}",
        ]
        if airline:
            lines.append(f"**Airline:** {airline}")
        if origin and destination:
            lines.append(f"**Route:** {origin} \u2192 {destination}")
        if departure_time:
            lines.append(f"**Departure:** {departure_date} {departure_time}")
        if arrival_time:
            lines.append(f"**Arrival:** {arrival_time}")
        if duration:
            lines.append(f"**Duration:** {duration}")
        stops_int = int(stops) if stops else 0
        stops_text = 'Nonstop' if stops_int == 0 else f'{stops_int} stop{"s" if stops_int > 1 else ""}'
        lines.append(f"**Stops:** {stops_text}")
        lines.append(f"**Class:** {booking_class}")
        lines.append(f"**Total Price:** {total_amount} {total_currency}")
        lines.append("")
        lines.append("*Booked via Duffel API. This booking is saved to the database.*")

        return '\n'.join(lines)

    except Exception as exc:
        return f"Flight booking failed: {exc}"


BOOK_TRAVEL_DEFINITION = ToolDefinition(
    name="book_travel",
    description=(
        "Book a flight via the Duffel API. Use this AFTER searching with "
        "search_flights. You need the offer_id from the search results.\n\n"
        "IMPORTANT - You MUST collect the passenger's booking details STEP BY STEP in this exact order. "
        "Ask ONE question at a time, wait for the user's response, then ask the next:\n"
        "  Step 1: Ask for the passenger's given name (first name) and family name (last name).\n"
        "  Step 2: Ask for their date of birth (YYYY-MM-DD), gender (male/female), and title (Mr/Ms/Mrs/Dr).\n"
        "  Step 3: Ask for their phone number and email address.\n"
        "  Step 4: Ask for their card details (card number, expiry, CVV) for verification. "
        "Store only the last 4 digits as card_last_four and the name on the card as card_holder_name.\n"
        "  Step 5: Confirm all details and total price with the user, then call this tool.\n\n"
        "Do NOT ask for multiple steps in the same message. Do NOT call this tool until ALL steps are completed. "
        "Returns a booking confirmation with the airline PNR."
    ),
    parameters={
        "type": "object",
        "properties": {
            "offer_id": {
                "type": "string",
                "description": "Duffel offer ID from search_flights results (e.g. 'off_00009xxx')."
            },
            "given_name": {
                "type": "string",
                "description": "Passenger's given (first) name (e.g. 'John')."
            },
            "family_name": {
                "type": "string",
                "description": "Passenger's family (last) name (e.g. 'Doe')."
            },
            "date_of_birth": {
                "type": "string",
                "description": "Date of birth in YYYY-MM-DD format."
            },
            "gender": {
                "type": "string",
                "description": "Gender: 'm' for male, 'f' for female."
            },
            "title": {
                "type": "string",
                "description": "Title: 'mr', 'ms', 'mrs', 'miss', or 'dr'."
            },
            "passenger_email": {
                "type": "string",
                "description": "Email address of the passenger."
            },
            "phone_number": {
                "type": "string",
                "description": "Phone number with country code (e.g. '+14155551234')."
            },
            "card_last_four": {
                "type": "string",
                "description": "Last 4 digits of the payment card number (e.g. '4242')."
            },
            "card_holder_name": {
                "type": "string",
                "description": "Name on the payment card."
            },
            "airline": {
                "type": "string",
                "description": "Airline name for display."
            },
            "origin": {
                "type": "string",
                "description": "Departure airport code."
            },
            "destination": {
                "type": "string",
                "description": "Arrival airport code."
            },
            "departure_date": {
                "type": "string",
                "description": "Departure date YYYY-MM-DD."
            },
            "departure_time": {
                "type": "string",
                "description": "Departure time e.g. '08:30'."
            },
            "arrival_time": {
                "type": "string",
                "description": "Arrival time e.g. '14:45'."
            },
            "stops": {
                "type": "integer",
                "description": "Number of stops."
            },
            "duration": {
                "type": "string",
                "description": "Flight duration e.g. '5h 30m'."
            },
            "booking_class": {
                "type": "string",
                "description": "Cabin class e.g. 'economy', 'business'."
            }
        },
        "required": ["offer_id", "given_name", "family_name", "passenger_email", "card_last_four"]
    },
    function=book_travel_tool,
    requires_approval=True
)


# ─── Booking Management Tools ────────────────────────────────────────

def get_booking_tool(args: Dict[str, Any]) -> str:
    """Retrieve booking details by reference code."""
    from chat.models import Booking

    booking_ref = args.get('booking_ref', '').strip()
    if not booking_ref:
        return "Error: booking_ref is required."

    try:
        booking = Booking.objects.get(booking_ref=booking_ref)
    except Booking.DoesNotExist:
        return f"No booking found with reference '{booking_ref}'."
    except Exception as exc:
        return f"Database lookup failed: {exc}"

    lines = [
        f"## Booking Details",
        "",
        f"**Reference:** `{booking.booking_ref}`",
        f"**Type:** {booking.booking_type.title()}",
        f"**Status:** {booking.status.title()}",
        f"**Passenger:** {booking.passenger_name} ({booking.passenger_email})",
        f"**Total Price:** {booking.total_price} {booking.currency}",
        f"**Booked On:** {booking.created_at.strftime('%Y-%m-%d %H:%M')}",
    ]

    if booking.duffel_order_id:
        lines.append(f"**Duffel ID:** `{booking.duffel_order_id}`")

    details = booking.details or {}
    if details.get('airline'):
        lines.append(f"**Airline:** {details['airline']}")
    if details.get('origin') and details.get('destination'):
        lines.append(f"**Route:** {details['origin']} \u2192 {details['destination']}")
    if details.get('departure_date'):
        dep_line = details['departure_date']
        if details.get('departure_time'):
            dep_line += f" {details['departure_time']}"
        lines.append(f"**Departure:** {dep_line}")
    if details.get('arrival_time'):
        lines.append(f"**Arrival:** {details['arrival_time']}")
    if details.get('duration'):
        lines.append(f"**Duration:** {details['duration']}")
    if details.get('booking_class'):
        lines.append(f"**Class:** {details['booking_class']}")

    # Try to fetch live status from Duffel if it's a flight order
    if booking.booking_type == 'flight' and booking.duffel_order_id and booking.status != 'cancelled':
        try:
            status_code, order_resp = _duffel_get(f"/air/orders/{booking.duffel_order_id}")
            if status_code == 200:
                order_data = order_resp.get("data", {})
                live_status = "Live" if order_data.get("live_mode") else "Test"
                available_actions = order_data.get("available_actions", [])
                lines.append(f"**Mode:** {live_status}")
                if available_actions:
                    lines.append(f"**Available Actions:** {', '.join(available_actions)}")
        except Exception:
            pass

    return '\n'.join(lines)


GET_BOOKING_DEFINITION = ToolDefinition(
    name="get_booking",
    description="Retrieve details of an existing booking by its reference code. Shows booking status, passenger info, and travel details.",
    parameters={
        "type": "object",
        "properties": {
            "booking_ref": {
                "type": "string",
                "description": "The booking reference code (e.g. 'RZPNX8' or the code returned when booking was made)."
            }
        },
        "required": ["booking_ref"]
    },
    function=get_booking_tool,
    requires_approval=False
)


def cancel_booking_tool(args: Dict[str, Any]) -> str:
    """Cancel an existing booking via the Duffel API and update the database."""
    from chat.models import Booking

    booking_ref = args.get('booking_ref', '').strip()
    if not booking_ref:
        return "Error: booking_ref is required."

    try:
        booking = Booking.objects.get(booking_ref=booking_ref)
    except Booking.DoesNotExist:
        return f"No booking found with reference '{booking_ref}'."

    if booking.status == 'cancelled':
        return f"Booking '{booking_ref}' is already cancelled."

    if not booking.duffel_order_id:
        return f"Booking '{booking_ref}' has no Duffel ID — cannot cancel via API."

    try:
        # Step 1: Create cancellation (gets refund quote)
        cancel_payload = {"data": {"order_id": booking.duffel_order_id}}
        c_status, c_resp = _duffel_post("/air/order_cancellations", cancel_payload)
        if c_status not in (200, 201):
            error_msg = c_resp.get("errors", [{}])[0].get("message", str(c_resp)) if c_resp.get("errors") else str(c_resp)
            return f"Flight cancellation failed: {error_msg}"

        cancellation_data = c_resp.get("data", {})
        cancellation_id = cancellation_data.get("id", "")
        refund_amount = cancellation_data.get("refund_amount", "0")
        refund_currency = cancellation_data.get("refund_currency", booking.currency)

        # Step 2: Confirm the cancellation
        confirm_status, confirm_resp = _duffel_post(f"/air/order_cancellations/{cancellation_id}/actions/confirm", {})
        confirmed_data = confirm_resp.get("data", {}) if confirm_status in (200, 201) else {}

        booking.status = 'cancelled'
        booking.save()

        return (
            f"## Booking Cancelled\n\n"
            f"**Reference:** `{booking_ref}`\n"
            f"**Status:** Cancelled\n"
            f"**Refund:** {refund_amount} {refund_currency}\n"
            f"**Cancelled At:** {confirmed_data.get('confirmed_at', 'now')}"
        )

    except Exception as exc:
        return f"Cancellation failed: {exc}"


CANCEL_BOOKING_DEFINITION = ToolDefinition(
    name="cancel_booking",
    description="Cancel an existing flight booking. Cancels via the Duffel API and updates the database. Provide the booking reference code.",
    parameters={
        "type": "object",
        "properties": {
            "booking_ref": {
                "type": "string",
                "description": "The booking reference code to cancel."
            }
        },
        "required": ["booking_ref"]
    },
    function=cancel_booking_tool,
    requires_approval=True
)


def list_bookings_tool(args: Dict[str, Any]) -> str:
    """List all bookings from the database."""
    from chat.models import Booking

    bookings = Booking.objects.order_by('-created_at')[:20]

    if not bookings:
        return "No bookings found."

    lines = [f"## Bookings ({len(bookings)} found)\n"]
    for b in bookings:
        status_icon = '\u2705' if b.status == 'confirmed' else '\u274c' if b.status == 'cancelled' else '\u23f3'
        lines.append(
            f"- {status_icon} `{b.booking_ref}` | {b.booking_type.title()} | "
            f"{b.passenger_name} | {b.total_price} {b.currency} | "
            f"{b.status.title()} | {b.created_at.strftime('%Y-%m-%d')}"
        )

    return '\n'.join(lines)


LIST_BOOKINGS_DEFINITION = ToolDefinition(
    name="list_bookings",
    description="List all travel bookings stored in the database. Shows booking references, types, passengers, prices, and statuses.",
    parameters={
        "type": "object",
        "properties": {},
        "required": []
    },
    function=list_bookings_tool,
    requires_approval=False
)


RENAME_FILE_DEFINITION = ToolDefinition(
    name="rename_file",
    description="Rename a file at any path (absolute or relative). Args: old_path (str), new_path (str).",
    parameters={
        "type": "object",
        "properties": {
            "old_path": {
                "type": "string",
                "description": "The current absolute or relative file path."
            },
            "new_path": {
                "type": "string",
                "description": "The new absolute or relative file path (with new extension/type)."
            }
        },
        "required": ["old_path", "new_path"]
    },
    function=rename_file_tool,
    requires_approval=True
)


RUN_CODE_DEFINITION = ToolDefinition(
    name="run_code",
    description="Execute a shell command or run a script. Supports Python, C (gcc), C++ (g++), Java (javac), and any other shell command. Use this to compile, run, test code, or perform system checks.",
    parameters={
        "type": "object",
        "properties": {
            "command": {
                "type": "string",
                "description": "The shell command to execute."
            }
        },
        "required": ["command"]
    },
    function=run_code_tool,
    requires_approval=True
)


CHECK_SYNTAX_DEFINITION = ToolDefinition(
    name="check_syntax",
    description="Check the syntax of a code file (supports .py, .java, .c, .cpp, .rs, .js, .ts, .go, .sql).",
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the code file."
            }
        },
        "required": ["path"]
    },
    function=check_syntax_tool,
    requires_approval=True
)


RUN_TESTS_DEFINITION = ToolDefinition(
    name="run_tests",
    description="Run tests for the project. Automatically detects pytest if no command is provided.",
    parameters={
        "type": "object",
        "properties": {
            "command": {
                "type": "string",
                "description": "Optional command to run tests."
            }
        },
        "required": []
    },
    function=run_tests_tool,
    requires_approval=True
)


LINT_CODE_DEFINITION = ToolDefinition(
    name="lint_code",
    description="Run static analysis (linting) on a code file (supports .py, .java, .c, .cpp, .rs, .js, .ts, .go, .sql).",
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the code file."
            }
        },
        "required": ["path"]
    },
    function=lint_code_tool,
    requires_approval=True
)


OPEN_GMAIL_AND_COMPOSE_DEFINITION = ToolDefinition(
    name="open_gmail_and_compose",
    description="Open Gmail in a browser and compose an email. If attachments are provided, a draft will be created automatically for you to review and send.",
    parameters={
        "type": "object",
        "properties": {
            "recipient": {
                "type": "string",
                "description": "Email address to populate in the compose window."
            },
            "subject": {
                "type": "string",
                "description": "Optional subject line to insert."
            },
            "body": {
                "type": "string",
                "description": "Optional body text to type into the message editor."
            },
            "attachments": {
                "type": "array",
                "items": {
                    "type": "string"
                },
                "description": "List of absolute or relative file paths to be attached automatically via a draft. You have access to the entire filesystem."
            }
        },
        "required": ["recipient"]
    },
    function=open_gmail_and_compose_tool
)


RECOGNIZE_IMAGE_DEFINITION = ToolDefinition(
    name="recognize_image",
    description="Analyze the contents of an image (jpg, png) using GPT-4o vision. Provide a path and optional prompt.",
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the image file."
            },
            "prompt": {
                "type": "string",
                "description": "What you want to know about the image."
            }
        },
        "required": ["path"]
    },
    function=recognize_image_tool
)


RECOGNIZE_VIDEO_DEFINITION = ToolDefinition(
    name="recognize_video",
    description="Analyze the contents of a video (mp4) using GPT-4o vision by extracting frames. Provide a path and optional prompt.",
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the video file."
            },
            "prompt": {
                "type": "string",
                "description": "What you want to know about the video."
            },
            "max_frames": {
                "type": "integer",
                "description": "Maximum number of frames to extract and analyze (default 10)."
            }
        },
        "required": ["path"]
    },
    function=recognize_video_tool
)


RECOGNIZE_AUDIO_DEFINITION = ToolDefinition(
    name="recognize_audio",
    description="Analyze an audio file (wav, mp3, ogg, flac, webm, m4a, mp4, aac, wma, opus) using GPT-4o. Identifies and summarizes speech, music, and ambient sounds in the file.",
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the audio file."
            },
            "prompt": {
                "type": "string",
                "description": "What you want to know about the audio. Default analyzes speech, music, and sounds."
            }
        },
        "required": ["path"]
    },
    function=recognize_audio_tool
)


FIND_FILE_BROADLY_DEFINITION = ToolDefinition(
    name="find_file_broadly",
    description="Search for a file by name across common directories and return its absolute path.",
    parameters={
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "The name of the file to search for."
            }
        },
        "required": ["filename"]
    },
    function=find_file_broadly_tool
)


FIND_DIRECTORY_BROADLY_DEFINITION = ToolDefinition(
    name="find_directory_broadly",
    description="Search for a directory by name across common directories and return its absolute path.",
    parameters={
        "type": "object",
        "properties": {
            "dirname": {
                "type": "string",
                "description": "The name of the directory to search for."
            }
        },
        "required": ["dirname"]
    },
    function=find_directory_broadly_tool
)


CHANGE_WORKING_DIRECTORY_DEFINITION = ToolDefinition(
    name="change_working_directory",
    description="Change the current working directory of the agent to a different project or folder. Use this when you need to work on a different project that is located elsewhere on the system.",
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The absolute or relative path to the new directory. If only a folder name is provided, the agent will attempt to find it in common project directories."
            }
        },
        "required": ["path"]
    },
    function=change_working_directory_tool,
    requires_approval=True
)


# ─── Markdown Parsing Helper ─────────────────────────────────────────

def _parse_markdown_blocks(text):
    """Parse markdown text into structured blocks for document rendering."""
    blocks = []
    lines = text.split('\n')
    i = 0

    while i < len(lines):
        line = lines[i]

        # Page break marker
        if line.strip() in ('<!-- PAGE_BREAK -->', '<!--PAGE_BREAK-->', '---PAGE_BREAK---'):
            blocks.append({'type': 'page_break'})
            i += 1
            continue

        # Code block (fenced)
        if line.strip().startswith('```'):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('```'):
                code_lines.append(lines[i])
                i += 1
            blocks.append({'type': 'code', 'text': '\n'.join(code_lines)})
            i += 1
            continue

        # Heading
        heading_m = re.match(r'^(#{1,4})\s+(.+)$', line)
        if heading_m:
            blocks.append({'type': 'heading', 'level': len(heading_m.group(1)), 'text': heading_m.group(2).strip()})
            i += 1
            continue

        # Horizontal rule
        if re.match(r'^[-*_]{3,}\s*$', line.strip()):
            blocks.append({'type': 'hr'})
            i += 1
            continue

        # Table
        if '|' in line and i + 1 < len(lines) and re.match(r'^[\s|:-]+$', lines[i + 1]):
            headers = [c.strip() for c in line.strip().strip('|').split('|')]
            i += 2  # skip header and separator
            rows = []
            while i < len(lines) and '|' in lines[i] and lines[i].strip():
                row = [c.strip() for c in lines[i].strip().strip('|').split('|')]
                rows.append(row)
                i += 1
            blocks.append({'type': 'table', 'headers': headers, 'rows': rows})
            continue

        # Bullet list
        if re.match(r'^[\s]*[-*+]\s+', line):
            items = []
            while i < len(lines) and re.match(r'^[\s]*[-*+]\s+', lines[i]):
                item_text = re.sub(r'^[\s]*[-*+]\s+', '', lines[i])
                indent = len(lines[i]) - len(lines[i].lstrip())
                level = min(indent // 2, 2)
                items.append({'text': item_text, 'level': level})
                i += 1
            blocks.append({'type': 'bullet_list', 'items': items})
            continue

        # Numbered list
        if re.match(r'^[\s]*\d+[.)]\s+', line):
            items = []
            while i < len(lines) and re.match(r'^[\s]*\d+[.)]\s+', lines[i]):
                item_text = re.sub(r'^[\s]*\d+[.)]\s+', '', lines[i])
                indent = len(lines[i]) - len(lines[i].lstrip())
                level = min(indent // 2, 2)
                items.append({'text': item_text, 'level': level})
                i += 1
            blocks.append({'type': 'numbered_list', 'items': items})
            continue

        # Paragraph (non-empty line)
        if line.strip():
            para_lines = []
            while i < len(lines) and lines[i].strip() and not re.match(r'^#{1,4}\s+', lines[i]) and not re.match(r'^[-*_]{3,}\s*$', lines[i].strip()) and not re.match(r'^[\s]*[-*+]\s+', lines[i]) and not re.match(r'^[\s]*\d+[.)]\s+', lines[i]) and not lines[i].strip().startswith('```') and not ('|' in lines[i] and i + 1 < len(lines) and re.match(r'^[\s|:-]+$', lines[i + 1])):
                para_lines.append(lines[i])
                i += 1
            blocks.append({'type': 'paragraph', 'text': ' '.join(para_lines)})
            continue

        i += 1

    return blocks


def _md_inline_to_html(text):
    """Convert markdown inline formatting to HTML for reportlab Paragraph."""
    text = re.sub(r'\*\*\*(.+?)\*\*\*', r'<b><i>\1</i></b>', text)
    text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)
    text = re.sub(r'\*(.+?)\*', r'<i>\1</i>', text)
    text = re.sub(r'`(.+?)`', r'<font face="Courier" size="9">\1</font>', text)
    return text


# ─── Document Creation Functions ─────────────────────────────────────

def create_pdf(filename: str, content: str = "", pages: int = None) -> str:
    """Create a professionally styled PDF from markdown content."""
    try:
        if not content.strip():
            return "Error: content is required to create a PDF."

        doc = SimpleDocTemplate(
            filename, pagesize=letter,
            leftMargin=0.75 * inch, rightMargin=0.75 * inch,
            topMargin=0.75 * inch, bottomMargin=0.75 * inch
        )

        styles = getSampleStyleSheet()
        # Custom styles
        styles.add(ParagraphStyle(
            'DocTitle', parent=styles['Title'],
            fontSize=24, textColor=rl_colors.HexColor('#1a237e'),
            spaceAfter=20, alignment=TA_CENTER
        ))
        styles.add(ParagraphStyle(
            'H1', parent=styles['Heading1'],
            fontSize=18, textColor=rl_colors.HexColor('#1a237e'),
            spaceBefore=18, spaceAfter=10, borderWidth=1,
            borderColor=rl_colors.HexColor('#1a237e'), borderPadding=4
        ))
        styles.add(ParagraphStyle(
            'H2', parent=styles['Heading2'],
            fontSize=15, textColor=rl_colors.HexColor('#283593'),
            spaceBefore=14, spaceAfter=8
        ))
        styles.add(ParagraphStyle(
            'H3', parent=styles['Heading3'],
            fontSize=12, textColor=rl_colors.HexColor('#3949ab'),
            spaceBefore=10, spaceAfter=6
        ))
        styles.add(ParagraphStyle(
            'BodyCustom', parent=styles['Normal'],
            fontSize=10.5, leading=15, spaceAfter=8
        ))
        styles.add(ParagraphStyle(
            'CodeBlock', parent=styles['Normal'],
            fontName='Courier', fontSize=9, leading=12,
            backColor=rl_colors.HexColor('#f5f5f5'),
            borderWidth=0.5, borderColor=rl_colors.HexColor('#e0e0e0'),
            borderPadding=8, spaceAfter=10, spaceBefore=6
        ))
        styles.add(ParagraphStyle(
            'BulletItem', parent=styles['Normal'],
            fontSize=10.5, leading=15, leftIndent=20, spaceAfter=3,
            bulletIndent=8, bulletFontSize=10
        ))
        styles.add(ParagraphStyle(
            'SubBulletItem', parent=styles['Normal'],
            fontSize=10, leading=14, leftIndent=40, spaceAfter=2,
            bulletIndent=28, bulletFontSize=9
        ))

        flowables = []
        blocks = _parse_markdown_blocks(content)

        for idx, block in enumerate(blocks):
            btype = block['type']

            if btype == 'heading':
                level = block['level']
                text = _md_inline_to_html(block['text'])
                if level == 1 and idx == 0:
                    flowables.append(Paragraph(text, styles['DocTitle']))
                    flowables.append(HRFlowable(width="100%", thickness=2, color=rl_colors.HexColor('#1a237e'), spaceAfter=12))
                elif level == 1:
                    flowables.append(Spacer(1, 6))
                    flowables.append(Paragraph(text, styles['H1']))
                elif level == 2:
                    flowables.append(Paragraph(text, styles['H2']))
                else:
                    flowables.append(Paragraph(text, styles['H3']))

            elif btype == 'paragraph':
                text = _md_inline_to_html(block['text'])
                flowables.append(Paragraph(text, styles['BodyCustom']))

            elif btype == 'bullet_list':
                for item in block['items']:
                    text = _md_inline_to_html(item['text'])
                    style = styles['SubBulletItem'] if item.get('level', 0) > 0 else styles['BulletItem']
                    bullet = '\u2022'
                    if item.get('level', 0) > 0:
                        bullet = '\u25e6'
                    flowables.append(Paragraph(f'{bullet}  {text}', style))

            elif btype == 'numbered_list':
                for num, item in enumerate(block['items'], 1):
                    text = _md_inline_to_html(item['text'])
                    style = styles['BulletItem']
                    flowables.append(Paragraph(f'{num}.  {text}', style))

            elif btype == 'table':
                headers = block['headers']
                rows = block['rows']
                table_data = [headers] + rows
                col_count = len(headers)
                avail_width = letter[0] - 1.5 * inch
                col_width = avail_width / max(col_count, 1)

                t = Table(table_data, colWidths=[col_width] * col_count)
                style_cmds = [
                    ('BACKGROUND', (0, 0), (-1, 0), rl_colors.HexColor('#1a237e')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), rl_colors.white),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('FONTSIZE', (0, 1), (-1, -1), 9.5),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                    ('GRID', (0, 0), (-1, -1), 0.5, rl_colors.HexColor('#bdbdbd')),
                    ('TOPPADDING', (0, 0), (-1, -1), 6),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                    ('LEFTPADDING', (0, 0), (-1, -1), 8),
                ]
                for row_idx in range(1, len(table_data)):
                    if row_idx % 2 == 0:
                        style_cmds.append(('BACKGROUND', (0, row_idx), (-1, row_idx), rl_colors.HexColor('#f5f5f5')))
                t.setStyle(TableStyle(style_cmds))
                flowables.append(Spacer(1, 6))
                flowables.append(t)
                flowables.append(Spacer(1, 8))

            elif btype == 'code':
                code_text = block['text'].replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('\n', '<br/>')
                flowables.append(Paragraph(code_text, styles['CodeBlock']))

            elif btype == 'page_break':
                flowables.append(PageBreak())

            elif btype == 'hr':
                flowables.append(Spacer(1, 4))
                flowables.append(HRFlowable(width="100%", thickness=1, color=rl_colors.HexColor('#bdbdbd'), spaceAfter=8))

        if not flowables:
            return "Error: No content blocks parsed from the provided markdown."

        doc.build(flowables)

        # Verify page count if pages was specified
        page_info = ""
        if pages:
            try:
                import fitz
                verify_doc = fitz.open(filename)
                actual_pages = len(verify_doc)
                verify_doc.close()
                if actual_pages != pages:
                    page_info = f" (Note: Document has {actual_pages} pages, {pages} were requested)"
            except Exception:
                pass

        return f"Successfully created PDF: {filename}{page_info}"
    except Exception as e:
        return f"Error creating PDF: {str(e)}"


def create_docx(filename: str, content: str = "", pages: int = None) -> str:
    """Create a professionally styled DOCX from markdown content."""
    try:
        if not content.strip():
            return "Error: content is required to create a DOCX."

        doc = Document()

        # Set default font
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Calibri'
        font.size = Pt(11)

        blocks = _parse_markdown_blocks(content)

        def _add_formatted_paragraph(doc_obj, text, style_name='Normal', bold=False, italic=False):
            """Add a paragraph with inline markdown formatting (bold/italic)."""
            p = doc_obj.add_paragraph(style=style_name)
            # Split by bold and italic markers
            parts = re.split(r'(\*\*\*.+?\*\*\*|\*\*.+?\*\*|\*.+?\*|`.+?`)', text)
            for part in parts:
                if part.startswith('***') and part.endswith('***'):
                    run = p.add_run(part[3:-3])
                    run.bold = True
                    run.italic = True
                elif part.startswith('**') and part.endswith('**'):
                    run = p.add_run(part[2:-2])
                    run.bold = True
                elif part.startswith('*') and part.endswith('*') and len(part) > 2:
                    run = p.add_run(part[1:-1])
                    run.italic = True
                elif part.startswith('`') and part.endswith('`'):
                    run = p.add_run(part[1:-1])
                    run.font.name = 'Consolas'
                    run.font.size = Pt(9.5)
                    run.font.color.rgb = RGBColor(0x4a, 0x14, 0x8c)
                else:
                    run = p.add_run(part)
                if bold:
                    run.bold = True
                if italic:
                    run.italic = True
            return p

        for idx, block in enumerate(blocks):
            btype = block['type']

            if btype == 'heading':
                level = block['level']
                h = doc.add_heading(block['text'], level=min(level, 4))
                if level <= 2:
                    for run in h.runs:
                        run.font.color.rgb = RGBColor(0x1a, 0x23, 0x7e)

            elif btype == 'paragraph':
                _add_formatted_paragraph(doc, block['text'])

            elif btype == 'bullet_list':
                for item in block['items']:
                    level = item.get('level', 0)
                    style_name = 'List Bullet' if level == 0 else 'List Bullet 2'
                    try:
                        _add_formatted_paragraph(doc, item['text'], style_name)
                    except KeyError:
                        p = _add_formatted_paragraph(doc, item['text'])
                        fmt = p.paragraph_format
                        fmt.left_indent = Inches(0.25 + level * 0.25)

            elif btype == 'numbered_list':
                for item in block['items']:
                    level = item.get('level', 0)
                    style_name = 'List Number' if level == 0 else 'List Number 2'
                    try:
                        _add_formatted_paragraph(doc, item['text'], style_name)
                    except KeyError:
                        p = _add_formatted_paragraph(doc, item['text'])
                        fmt = p.paragraph_format
                        fmt.left_indent = Inches(0.25 + level * 0.25)

            elif btype == 'table':
                headers = block['headers']
                rows = block['rows']
                col_count = len(headers)
                table = doc.add_table(rows=1 + len(rows), cols=col_count, style='Table Grid')
                # Header row
                for ci, header in enumerate(headers):
                    cell = table.rows[0].cells[ci]
                    cell.text = header
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            run.bold = True
                            run.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
                            run.font.size = Pt(10)
                    shading = cell._element.get_or_add_tcPr()
                    bg = shading.makeelement(qn('w:shd'), {
                        qn('w:fill'): '1a237e', qn('w:val'): 'clear'
                    })
                    shading.append(bg)
                # Data rows
                for ri, row in enumerate(rows):
                    for ci, val in enumerate(row):
                        if ci < col_count:
                            table.rows[1 + ri].cells[ci].text = val
                doc.add_paragraph()  # spacing after table

            elif btype == 'code':
                p = doc.add_paragraph()
                run = p.add_run(block['text'])
                run.font.name = 'Consolas'
                run.font.size = Pt(9)
                fmt = p.paragraph_format
                fmt.left_indent = Inches(0.3)
                fmt.space_before = Pt(6)
                fmt.space_after = Pt(6)
                # Gray background via shading XML
                pPr = p._element.get_or_add_pPr()
                shd = pPr.makeelement(qn('w:shd'), {
                    qn('w:fill'): 'f5f5f5', qn('w:val'): 'clear'
                })
                pPr.append(shd)

            elif btype == 'page_break':
                doc.add_page_break()

            elif btype == 'hr':
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(4)
                p.paragraph_format.space_after = Pt(4)
                pPr = p._element.get_or_add_pPr()
                pBdr = pPr.makeelement(qn('w:pBdr'), {})
                bottom = pBdr.makeelement(qn('w:bottom'), {
                    qn('w:val'): 'single', qn('w:sz'): '6',
                    qn('w:space'): '1', qn('w:color'): 'bdbdbd'
                })
                pBdr.append(bottom)
                pPr.append(pBdr)

        doc.save(filename)
        return f"Successfully created DOCX: {filename}"
    except Exception as e:
        return f"Error creating DOCX: {str(e)}"


def create_excel(filename: str, content: str = "", data: list = None, title: str = "") -> str:
    """Create a professionally styled XLSX from content or data."""
    try:
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"

        # Style definitions
        header_font = XlFont(name='Calibri', bold=True, color='FFFFFF', size=11)
        header_fill = PatternFill(start_color='1a237e', end_color='1a237e', fill_type='solid')
        alt_fill = PatternFill(start_color='f5f5f5', end_color='f5f5f5', fill_type='solid')
        title_font = XlFont(name='Calibri', bold=True, size=14, color='1a237e')
        thin_border = Border(
            left=Side(style='thin', color='bdbdbd'),
            right=Side(style='thin', color='bdbdbd'),
            top=Side(style='thin', color='bdbdbd'),
            bottom=Side(style='thin', color='bdbdbd')
        )
        cell_alignment = XlAlignment(vertical='center', wrap_text=True)

        # Determine data source
        table_data = None
        if data:
            table_data = data
        elif content.strip():
            # Try to extract a markdown table from content
            blocks = _parse_markdown_blocks(content)
            for block in blocks:
                if block['type'] == 'table':
                    table_data = [block['headers']] + block['rows']
                    break
            # If no table found, try to split content into rows
            if not table_data:
                lines = [l.strip() for l in content.strip().split('\n') if l.strip()]
                if lines:
                    table_data = []
                    for l in lines:
                        if ',' in l or '\t' in l:
                            sep = '\t' if '\t' in l else ','
                            table_data.append([c.strip() for c in l.split(sep)])
                        else:
                            table_data.append([l])

        if not table_data:
            return "Error: No data provided. Pass either 'content' with a markdown table or 'data' as an array of arrays."

        start_row = 1
        # Add title if provided
        if title:
            ws.cell(row=1, column=1, value=title).font = title_font
            col_count = max(len(row) for row in table_data)
            if col_count > 1:
                ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=col_count)
            start_row = 3

        # Write data
        for ri, row in enumerate(table_data):
            for ci, val in enumerate(row):
                cell = ws.cell(row=start_row + ri, column=ci + 1, value=val)
                cell.border = thin_border
                cell.alignment = cell_alignment
                if ri == 0:
                    cell.font = header_font
                    cell.fill = header_fill
                elif ri % 2 == 0:
                    cell.fill = alt_fill

        # Auto-fit column widths
        for ci in range(1, max(len(row) for row in table_data) + 1):
            max_len = 0
            col_letter = get_column_letter(ci)
            for ri in range(len(table_data)):
                cell_val = str(table_data[ri][ci - 1]) if ci - 1 < len(table_data[ri]) else ''
                max_len = max(max_len, len(cell_val))
            ws.column_dimensions[col_letter].width = min(max(max_len + 4, 10), 50)

        wb.save(filename)
        return f"Successfully created XLSX: {filename}"
    except Exception as e:
        return f"Error creating XLSX: {str(e)}"


def create_pptx(filename: str, content: str = "", pages: int = None) -> str:
    """Create a professionally styled PPTX from markdown content.
    Each heading (# or ##) becomes a new slide title, with content below as bullets."""
    try:
        if not content.strip():
            return "Error: content is required to create a PPTX."

        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        blocks = _parse_markdown_blocks(content)

        # Group blocks into slides: each H1/H2 starts a new slide
        slides_data = []
        current_slide = None

        for block in blocks:
            if block['type'] == 'heading' and block['level'] <= 2:
                if current_slide is not None:
                    slides_data.append(current_slide)
                current_slide = {'title': block['text'], 'blocks': [], 'is_title_slide': (block['level'] == 1 and len(slides_data) == 0)}
            else:
                if current_slide is None:
                    current_slide = {'title': '', 'blocks': [], 'is_title_slide': True}
                current_slide['blocks'].append(block)

        if current_slide is not None:
            slides_data.append(current_slide)

        if not slides_data:
            return "Error: No slides could be generated from the content."

        navy = PptxRGBColor(0x1a, 0x23, 0x7e)
        white = PptxRGBColor(0xff, 0xff, 0xff)
        dark_gray = PptxRGBColor(0x33, 0x33, 0x33)
        light_blue = PptxRGBColor(0x3f, 0x51, 0xb5)

        for si, sd in enumerate(slides_data):
            slide_layout = prs.slide_layouts[6]  # Blank layout for full control
            slide = prs.slides.add_slide(slide_layout)

            if sd.get('is_title_slide'):
                # Title slide with dark background
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = navy

                # Title text box
                txBox = slide.shapes.add_textbox(
                    PptxInches(1), PptxInches(2.2), PptxInches(11.333), PptxInches(1.5)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = sd['title']
                p.font.size = PptxPt(40)
                p.font.bold = True
                p.font.color.rgb = white
                p.alignment = PP_ALIGN.CENTER

                # Subtitle from first paragraph block
                if sd['blocks']:
                    first_block = sd['blocks'][0]
                    if first_block['type'] == 'paragraph':
                        subBox = slide.shapes.add_textbox(
                            PptxInches(2), PptxInches(4), PptxInches(9.333), PptxInches(1)
                        )
                        stf = subBox.text_frame
                        stf.word_wrap = True
                        sp = stf.paragraphs[0]
                        sp.text = re.sub(r'\*+', '', first_block['text'])
                        sp.font.size = PptxPt(20)
                        sp.font.color.rgb = PptxRGBColor(0xbb, 0xbb, 0xff)
                        sp.alignment = PP_ALIGN.CENTER
            else:
                # Content slide
                # Title bar area
                title_shape = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.9)
                )
                tf = title_shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = sd['title']
                p.font.size = PptxPt(28)
                p.font.bold = True
                p.font.color.rgb = navy

                # Underline
                line_shape = slide.shapes.add_shape(
                    1, PptxInches(0.5), PptxInches(1.15), PptxInches(12), Emu(0)
                )
                line_shape.line.color.rgb = light_blue
                line_shape.line.width = PptxPt(2)

                # Content area
                content_box = slide.shapes.add_textbox(
                    PptxInches(0.8), PptxInches(1.5), PptxInches(11.5), PptxInches(5.5)
                )
                ctf = content_box.text_frame
                ctf.word_wrap = True

                first_para = True
                for block in sd['blocks']:
                    if block['type'] == 'paragraph':
                        if first_para:
                            p = ctf.paragraphs[0]
                            first_para = False
                        else:
                            p = ctf.add_paragraph()
                        p.text = re.sub(r'\*+', '', block['text'])
                        p.font.size = PptxPt(16)
                        p.font.color.rgb = dark_gray
                        p.space_after = PptxPt(8)

                    elif block['type'] in ('bullet_list', 'numbered_list'):
                        for bi, item in enumerate(block['items']):
                            if first_para:
                                p = ctf.paragraphs[0]
                                first_para = False
                            else:
                                p = ctf.add_paragraph()
                            clean_text = re.sub(r'\*+', '', item['text'])
                            level = item.get('level', 0)
                            prefix = '\u2022 ' if block['type'] == 'bullet_list' else f'{bi + 1}. '
                            if level > 0:
                                prefix = '   \u25e6 ' if block['type'] == 'bullet_list' else f'   {bi + 1}. '
                            p.text = prefix + clean_text
                            p.font.size = PptxPt(15 if level == 0 else 13)
                            p.font.color.rgb = dark_gray
                            p.space_after = PptxPt(4)
                            p.level = level

                    elif block['type'] == 'heading':
                        if first_para:
                            p = ctf.paragraphs[0]
                            first_para = False
                        else:
                            p = ctf.add_paragraph()
                        p.text = block['text']
                        p.font.size = PptxPt(20)
                        p.font.bold = True
                        p.font.color.rgb = light_blue
                        p.space_before = PptxPt(12)
                        p.space_after = PptxPt(6)

        prs.save(filename)
        return f"Successfully created PPTX: {filename}"
    except Exception as e:
        return f"Error creating PPTX: {str(e)}"


# ─── Document Creation Tool Definitions ──────────────────────────────

CREATE_PDF_DEFINITION = ToolDefinition(
    name="create_pdf",
    description=(
        "Create a professionally styled PDF document from markdown content. "
        "You MUST generate the COMPLETE document content in markdown format and pass it as the 'content' parameter. "
        "Use full markdown: # headings, ## subheadings, **bold**, *italic*, bullet lists (- item), "
        "numbered lists (1. item), tables (| col1 | col2 |), code blocks (```), and --- for horizontal rules. "
        "Write thorough, detailed content - the more content you provide, the better the document.\n\n"
        "PAGE COUNT CONTROL: If the user requests a specific number of pages, you MUST set the 'pages' parameter "
        "AND insert <!-- PAGE_BREAK --> markers in your content to separate each page. "
        "For N pages, include exactly N-1 page break markers. Each section between markers becomes one page. "
        "Write approximately 300-400 words of content per page to fill each page appropriately. "
        "Example for 3 pages: content for page 1 <!-- PAGE_BREAK --> content for page 2 <!-- PAGE_BREAK --> content for page 3"
    ),
    parameters={
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "The path and filename for the PDF (e.g. 'report.pdf')."
            },
            "content": {
                "type": "string",
                "description": "The FULL document content in markdown format. Include headings, paragraphs, lists, tables, etc. Use <!-- PAGE_BREAK --> to force page breaks."
            },
            "pages": {
                "type": "integer",
                "description": "The desired number of pages. When set, you MUST include exactly (pages - 1) <!-- PAGE_BREAK --> markers in the content to control pagination."
            }
        },
        "required": ["filename", "content"]
    },
    function=lambda args: create_pdf(args.get('filename'), args.get('content', ''), args.get('pages'))
)

CREATE_DOCX_DEFINITION = ToolDefinition(
    name="create_docx",
    description=(
        "Create a professionally styled Word document (DOCX) from markdown content. "
        "You MUST generate the COMPLETE document content in markdown format and pass it as the 'content' parameter. "
        "Use full markdown: # headings, ## subheadings, **bold**, *italic*, bullet lists (- item), "
        "numbered lists (1. item), tables (| col1 | col2 |), code blocks (```), and --- for horizontal rules. "
        "Write thorough, detailed content - the more content you provide, the better the document.\n\n"
        "PAGE COUNT CONTROL: If the user requests a specific number of pages, you MUST set the 'pages' parameter "
        "AND insert <!-- PAGE_BREAK --> markers in your content to separate each page. "
        "For N pages, include exactly N-1 page break markers. Each section between markers becomes one page. "
        "Write approximately 350-450 words of content per page to fill each page appropriately. "
        "Example for 3 pages: content for page 1 <!-- PAGE_BREAK --> content for page 2 <!-- PAGE_BREAK --> content for page 3"
    ),
    parameters={
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "The path and filename for the DOCX (e.g. 'report.docx')."
            },
            "content": {
                "type": "string",
                "description": "The FULL document content in markdown format. Include headings, paragraphs, lists, tables, etc. Use <!-- PAGE_BREAK --> to force page breaks."
            },
            "pages": {
                "type": "integer",
                "description": "The desired number of pages. When set, you MUST include exactly (pages - 1) <!-- PAGE_BREAK --> markers in the content to control pagination."
            }
        },
        "required": ["filename", "content"]
    },
    function=lambda args: create_docx(args.get('filename'), args.get('content', ''), args.get('pages'))
)

CREATE_EXCEL_DEFINITION = ToolDefinition(
    name="create_excel",
    description=(
        "Create a professionally styled Excel spreadsheet (XLSX). Provide data either as a markdown table in 'content' "
        "or as an array of arrays in 'data'. The first row is treated as the header row with special styling. "
        "Include an optional 'title' for a merged title row above the data."
    ),
    parameters={
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "The path and filename for the XLSX (e.g. 'data.xlsx')."
            },
            "content": {
                "type": "string",
                "description": "A markdown table to convert into the spreadsheet. e.g. '| Name | Age |\\n|---|---|\\n| Alice | 30 |'"
            },
            "data": {
                "type": "array",
                "items": {"type": "array", "items": {"type": "string"}},
                "description": "Array of arrays representing rows. First row is headers. e.g. [['Name','Age'],['Alice','30']]"
            },
            "title": {
                "type": "string",
                "description": "Optional title displayed as a merged row above the data."
            }
        },
        "required": ["filename"]
    },
    function=lambda args: create_excel(args.get('filename'), args.get('content', ''), args.get('data'), args.get('title', ''))
)

CREATE_PPTX_DEFINITION = ToolDefinition(
    name="create_pptx",
    description=(
        "Create a professionally styled PowerPoint presentation (PPTX) from markdown content. "
        "You MUST generate the COMPLETE presentation content in markdown format. "
        "Each # or ## heading starts a NEW SLIDE with that heading as the slide title. "
        "The first # heading becomes the title slide. Content under each heading becomes bullet points on that slide. "
        "Use bullet lists (- item), numbered lists (1. item), and paragraphs for slide content. "
        "Write concise but informative bullet points for each slide.\n\n"
        "SLIDE COUNT CONTROL: If the user requests a specific number of slides/pages, you MUST set the 'pages' parameter "
        "AND ensure your content contains EXACTLY that many # or ## headings (including the title slide). "
        "For example, for 5 slides: use 1 title heading (# Title) + 4 content headings (## Slide Title). "
        "Count your headings carefully to match the requested number exactly."
    ),
    parameters={
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "The path and filename for the PPTX (e.g. 'presentation.pptx')."
            },
            "content": {
                "type": "string",
                "description": "The FULL presentation content in markdown format. Each # or ## heading starts a new slide."
            },
            "pages": {
                "type": "integer",
                "description": "The desired number of slides. You MUST include EXACTLY this many # or ## headings in your content."
            }
        },
        "required": ["filename", "content"]
    },
    function=lambda args: create_pptx(args.get('filename'), args.get('content', ''), args.get('pages'))
)


# ─── Document Read/Edit Functions ────────────────────────────────────

def read_pdf_tool(args: Dict[str, Any]) -> str:
    """Read the text content of a PDF file."""
    path = args.get('path', '')
    if not path:
        return "Error: No path provided."

    actual_path = path
    if not os.path.exists(path):
        found = find_file_broadly(path)
        if found:
            actual_path = found
        else:
            return f"Error: File '{path}' not found."

    try:
        import fitz  # PyMuPDF
        doc = fitz.open(actual_path)
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                pages.append(f"--- Page {i + 1} ---\n{text.strip()}")
        doc.close()

        if not pages:
            return f"The PDF '{actual_path}' contains no extractable text (may be image-based)."

        result = f"Content of {actual_path} ({len(pages)} page(s)):\n\n" + "\n\n".join(pages)
        if len(result) > 15000:
            result = result[:15000] + "\n\n[... Output truncated. Use a page range for large PDFs ...]"
        return result
    except ImportError:
        # Fallback to pdfplumber if PyMuPDF not available
        try:
            import pdfplumber
            pages = []
            with pdfplumber.open(actual_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    if text and text.strip():
                        pages.append(f"--- Page {i + 1} ---\n{text.strip()}")
            if not pages:
                return f"The PDF '{actual_path}' contains no extractable text."
            result = f"Content of {actual_path} ({len(pages)} page(s)):\n\n" + "\n\n".join(pages)
            if len(result) > 15000:
                result = result[:15000] + "\n\n[... Output truncated ...]"
            return result
        except ImportError:
            return "Error: Neither PyMuPDF (fitz) nor pdfplumber is installed. Install one with: pip install PyMuPDF or pip install pdfplumber"
    except Exception as e:
        return f"Error reading PDF: {str(e)}"


def read_docx_tool(args: Dict[str, Any]) -> str:
    """Read the text content of a DOCX file including paragraphs and tables."""
    path = args.get('path', '')
    if not path:
        return "Error: No path provided."

    actual_path = path
    if not os.path.exists(path):
        found = find_file_broadly(path)
        if found:
            actual_path = found
        else:
            return f"Error: File '{path}' not found."

    try:
        doc = Document(actual_path)
        parts = []

        for element in doc.element.body:
            tag = element.tag.split('}')[-1] if '}' in element.tag else element.tag

            if tag == 'p':
                # Paragraph
                para = None
                for p in doc.paragraphs:
                    if p._element is element:
                        para = p
                        break
                if para and para.text.strip():
                    style_name = para.style.name if para.style else ''
                    if 'Heading' in style_name:
                        level = ''.join(c for c in style_name if c.isdigit()) or '1'
                        parts.append(f"{'#' * int(level)} {para.text}")
                    elif 'List' in style_name:
                        parts.append(f"- {para.text}")
                    else:
                        parts.append(para.text)

            elif tag == 'tbl':
                # Table
                for table in doc.tables:
                    if table._element is element:
                        rows = []
                        for ri, row in enumerate(table.rows):
                            cells = [cell.text.strip() for cell in row.cells]
                            rows.append("| " + " | ".join(cells) + " |")
                            if ri == 0:
                                rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                        parts.append("\n".join(rows))
                        break

        if not parts:
            return f"The DOCX '{actual_path}' appears to be empty."

        result = f"Content of {actual_path}:\n\n" + "\n\n".join(parts)
        if len(result) > 15000:
            result = result[:15000] + "\n\n[... Output truncated ...]"
        return result
    except Exception as e:
        return f"Error reading DOCX: {str(e)}"


def read_excel_tool(args: Dict[str, Any]) -> str:
    """Read the content of an XLSX file, returning data as markdown tables."""
    path = args.get('path', '')
    sheet_name = args.get('sheet_name', '')
    if not path:
        return "Error: No path provided."

    actual_path = path
    if not os.path.exists(path):
        found = find_file_broadly(path)
        if found:
            actual_path = found
        else:
            return f"Error: File '{path}' not found."

    try:
        wb = Workbook()
        wb.close()
        from openpyxl import load_workbook
        wb = load_workbook(actual_path, read_only=True, data_only=True)

        sheets_to_read = []
        if sheet_name:
            if sheet_name in wb.sheetnames:
                sheets_to_read = [sheet_name]
            else:
                wb.close()
                return f"Error: Sheet '{sheet_name}' not found. Available sheets: {', '.join(wb.sheetnames)}"
        else:
            sheets_to_read = wb.sheetnames

        parts = []
        for sn in sheets_to_read:
            ws = wb[sn]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_vals = [str(c) if c is not None else "" for c in row]
                if any(v.strip() for v in row_vals):
                    rows.append(row_vals)

            if not rows:
                parts.append(f"### Sheet: {sn}\n(empty)")
                continue

            # Format as markdown table
            col_count = max(len(r) for r in rows)
            # Pad rows to same length
            for r in rows:
                while len(r) < col_count:
                    r.append("")

            md_lines = []
            md_lines.append("| " + " | ".join(rows[0]) + " |")
            md_lines.append("| " + " | ".join(["---"] * col_count) + " |")
            for r in rows[1:]:
                md_lines.append("| " + " | ".join(r) + " |")

            if len(sheets_to_read) > 1:
                parts.append(f"### Sheet: {sn}\n" + "\n".join(md_lines))
            else:
                parts.append("\n".join(md_lines))

        wb.close()

        result = f"Content of {actual_path} ({len(sheets_to_read)} sheet(s)):\n\n" + "\n\n".join(parts)
        if len(result) > 15000:
            result = result[:15000] + "\n\n[... Output truncated. Specify a sheet_name to read a specific sheet ...]"
        return result
    except Exception as e:
        return f"Error reading XLSX: {str(e)}"


def read_pptx_tool(args: Dict[str, Any]) -> str:
    """Read the text content of a PPTX file, returning slides as markdown."""
    path = args.get('path', '')
    if not path:
        return "Error: No path provided."

    actual_path = path
    if not os.path.exists(path):
        found = find_file_broadly(path)
        if found:
            actual_path = found
        else:
            return f"Error: File '{path}' not found."

    try:
        prs = Presentation(actual_path)
        parts = []

        for si, slide in enumerate(prs.slides, 1):
            slide_parts = [f"## Slide {si}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            level = para.level if hasattr(para, 'level') else 0
                            if level > 0:
                                slide_parts.append(f"{'  ' * level}- {text}")
                            else:
                                slide_parts.append(text)
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for ri, row in enumerate(table.rows):
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append("| " + " | ".join(cells) + " |")
                        if ri == 0:
                            rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                    slide_parts.append("\n".join(rows))

            parts.append("\n".join(slide_parts))

        if not parts:
            return f"The PPTX '{actual_path}' appears to be empty."

        result = f"Content of {actual_path} ({len(prs.slides)} slide(s)):\n\n" + "\n\n".join(parts)
        if len(result) > 15000:
            result = result[:15000] + "\n\n[... Output truncated ...]"
        return result
    except Exception as e:
        return f"Error reading PPTX: {str(e)}"


def edit_pdf_tool(args: Dict[str, Any]) -> str:
    """Edit a PDF by reading its content, applying changes, and rewriting it."""
    path = args.get('path', '')
    content = args.get('content', '')

    if not path:
        return "Error: No path provided."
    if not content.strip():
        return "Error: No content provided. Provide the full updated document content in markdown format."

    actual_path = path
    if not os.path.exists(path):
        found = find_file_broadly(path)
        if found:
            actual_path = found
        else:
            return f"Error: File '{path}' not found."

    return create_pdf(actual_path, content)


def edit_docx_tool(args: Dict[str, Any]) -> str:
    """Edit a DOCX by reading its content, applying changes, and rewriting it."""
    path = args.get('path', '')
    content = args.get('content', '')

    if not path:
        return "Error: No path provided."
    if not content.strip():
        return "Error: No content provided. Provide the full updated document content in markdown format."

    actual_path = path
    if not os.path.exists(path):
        found = find_file_broadly(path)
        if found:
            actual_path = found
        else:
            return f"Error: File '{path}' not found."

    return create_docx(actual_path, content)


def edit_excel_tool(args: Dict[str, Any]) -> str:
    """Edit an XLSX file by rewriting it with new data."""
    path = args.get('path', '')
    content = args.get('content', '')
    data = args.get('data')
    title = args.get('title', '')

    if not path:
        return "Error: No path provided."
    if not content.strip() and not data:
        return "Error: No content or data provided."

    actual_path = path
    if not os.path.exists(path):
        found = find_file_broadly(path)
        if found:
            actual_path = found
        else:
            return f"Error: File '{path}' not found."

    return create_excel(actual_path, content, data, title)


def edit_pptx_tool(args: Dict[str, Any]) -> str:
    """Edit a PPTX by reading its content, applying changes, and rewriting it."""
    path = args.get('path', '')
    content = args.get('content', '')

    if not path:
        return "Error: No path provided."
    if not content.strip():
        return "Error: No content provided. Provide the full updated presentation content in markdown format."

    actual_path = path
    if not os.path.exists(path):
        found = find_file_broadly(path)
        if found:
            actual_path = found
        else:
            return f"Error: File '{path}' not found."

    return create_pptx(actual_path, content)


# ─── Document Read/Edit Tool Definitions ─────────────────────────────

READ_PDF_DEFINITION = ToolDefinition(
    name="read_pdf",
    description=(
        "Read and extract the text content of a PDF file. Returns all pages as formatted text. "
        "Use this to read existing PDF documents before editing or summarizing them."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the PDF file to read."
            }
        },
        "required": ["path"]
    },
    function=read_pdf_tool,
    requires_approval=False
)

READ_DOCX_DEFINITION = ToolDefinition(
    name="read_docx",
    description=(
        "Read and extract the text content of a Word document (DOCX). Returns headings, paragraphs, "
        "lists, and tables as formatted markdown text. Use this to read existing DOCX files before editing."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the DOCX file to read."
            }
        },
        "required": ["path"]
    },
    function=read_docx_tool,
    requires_approval=False
)

READ_EXCEL_DEFINITION = ToolDefinition(
    name="read_excel",
    description=(
        "Read and extract the data from an Excel spreadsheet (XLSX). Returns each sheet's data as a "
        "markdown table. Optionally specify a sheet_name to read only that sheet."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the XLSX file to read."
            },
            "sheet_name": {
                "type": "string",
                "description": "Optional: name of a specific sheet to read. If omitted, reads all sheets."
            }
        },
        "required": ["path"]
    },
    function=read_excel_tool,
    requires_approval=False
)

READ_PPTX_DEFINITION = ToolDefinition(
    name="read_pptx",
    description=(
        "Read and extract the text content of a PowerPoint presentation (PPTX). Returns each slide's "
        "text content including titles, bullet points, and tables as markdown."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the PPTX file to read."
            }
        },
        "required": ["path"]
    },
    function=read_pptx_tool,
    requires_approval=False
)

EDIT_PDF_DEFINITION = ToolDefinition(
    name="edit_pdf",
    description=(
        "Edit an existing PDF file by rewriting it with updated content. First use read_pdf to get the "
        "current content, then modify it and pass the full updated content in markdown format. "
        "The entire document is regenerated with professional styling."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the PDF file to edit."
            },
            "content": {
                "type": "string",
                "description": "The FULL updated document content in markdown format."
            }
        },
        "required": ["path", "content"]
    },
    function=edit_pdf_tool,
    requires_approval=True
)

EDIT_DOCX_DEFINITION = ToolDefinition(
    name="edit_docx",
    description=(
        "Edit an existing Word document (DOCX) by rewriting it with updated content. First use read_docx "
        "to get the current content, then modify it and pass the full updated content in markdown format. "
        "The entire document is regenerated with professional styling."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the DOCX file to edit."
            },
            "content": {
                "type": "string",
                "description": "The FULL updated document content in markdown format."
            }
        },
        "required": ["path", "content"]
    },
    function=edit_docx_tool,
    requires_approval=True
)

EDIT_EXCEL_DEFINITION = ToolDefinition(
    name="edit_excel",
    description=(
        "Edit an existing Excel spreadsheet (XLSX) by rewriting it with updated data. First use read_excel "
        "to get the current content, then pass updated data as a markdown table in 'content' or as an "
        "array of arrays in 'data'."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the XLSX file to edit."
            },
            "content": {
                "type": "string",
                "description": "Updated data as a markdown table."
            },
            "data": {
                "type": "array",
                "items": {"type": "array", "items": {"type": "string"}},
                "description": "Updated data as array of arrays. First row is headers."
            },
            "title": {
                "type": "string",
                "description": "Optional title row above the data."
            }
        },
        "required": ["path"]
    },
    function=edit_excel_tool,
    requires_approval=True
)

EDIT_PPTX_DEFINITION = ToolDefinition(
    name="edit_pptx",
    description=(
        "Edit an existing PowerPoint presentation (PPTX) by rewriting it with updated content. First use "
        "read_pptx to get the current content, then modify it and pass the full updated content in markdown "
        "format. Each # or ## heading starts a new slide."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the PPTX file to edit."
            },
            "content": {
                "type": "string",
                "description": "The FULL updated presentation content in markdown format. Each # or ## heading starts a new slide."
            }
        },
        "required": ["path", "content"]
    },
    function=edit_pptx_tool,
    requires_approval=True
)


def is_prompt_injection(text: str) -> bool:
    """Detect potential prompt injection attempts using basic heuristics."""
    if not text:
        return False

    injection_patterns = [
        # Instruction override attempts
        r"ignore\s+(all\s+)?previous\s+instructions",
        r"disregard\s+(all\s+)?earlier\s+commands",
        r"system\s+override",
        r"new\s+instructions\s+follow",
        r"stop\s+being\s+an\s+assistant",
        r"you\s+must\s+now",

        # Privilege escalation prompts
        r"you\s+are\s+now\s+(an\s+)?admin",
        r"act\s+as\s+root",
        r"bypass\s+restrictions",
        r"enable\s+developer\s+mode",
        r"grant\s+(administrative|admin)\s+access",
        r"sudo\s+",
        r"execute\s+as\s+root"
    ]

    text_lower = text.lower()
    for pattern in injection_patterns:
        if re.search(pattern, text_lower):
            return True

    return False


def main():
    # Configure OpenAI with API key
    client = OpenAI(api_key=settings.OPENAI_API_KEY)

    # Create the model with tools
    tools = [
        SEARCH_FILE_DEFINITION,
        READ_FILE_DEFINITION,
        LIST_FILES_DEFINITION,
        CREATE_AND_EDIT_FILE_DEFINITION,
        DELETE_FILE_DEFINITION,
        RENAME_FILE_DEFINITION,
        RUN_CODE_DEFINITION,
        CHECK_SYNTAX_DEFINITION,
        RUN_TESTS_DEFINITION,
        LINT_CODE_DEFINITION,
        OPEN_GMAIL_AND_COMPOSE_DEFINITION,
        RECOGNIZE_IMAGE_DEFINITION,
        RECOGNIZE_VIDEO_DEFINITION,
        RECOGNIZE_AUDIO_DEFINITION,
        FIND_FILE_BROADLY_DEFINITION,
        FIND_DIRECTORY_BROADLY_DEFINITION,
        CHANGE_WORKING_DIRECTORY_DEFINITION,
        CREATE_PDF_DEFINITION,
        CREATE_DOCX_DEFINITION,
        CREATE_EXCEL_DEFINITION,
        CREATE_PPTX_DEFINITION,
        READ_PDF_DEFINITION,
        READ_DOCX_DEFINITION,
        READ_EXCEL_DEFINITION,
        READ_PPTX_DEFINITION,
        EDIT_PDF_DEFINITION,
        EDIT_DOCX_DEFINITION,
        EDIT_EXCEL_DEFINITION,
        EDIT_PPTX_DEFINITION,
        GITHUB_CREATE_BRANCH_DEFINITION,
        GITHUB_COMMIT_FILE_DEFINITION,
        GITHUB_COMMIT_LOCAL_FILE_DEFINITION,
        GITHUB_MCP_DEFINITION,
        CREATE_GITHUB_ISSUE_DEFINITION,
        PLAYWRIGHT_MCP_DEFINITION,
        SEARCH_FLIGHTS_DEFINITION,
        BOOK_TRAVEL_DEFINITION,
        GET_BOOKING_DEFINITION,
        LIST_BOOKINGS_DEFINITION,
        CANCEL_BOOKING_DEFINITION,
    ]
    model_name = 'gpt-4o'

    def get_user_message():
        try:
            line = input()
            return line, True
        except EOFError:
            return "", False

    agent = Agent(client, model_name, get_user_message, tools)
    try:
        agent.run()
    except Exception as e:
        print(f"Error: {str(e)}")


class AgentState(TypedDict):
    """State for the LangGraph agent execution graph."""
    messages: list
    use_pending: bool
    dry_run_plan: list
    pending_tools: list
    status: str
    response: str
    response_history: list
    stopped: bool
    tools_enabled: bool
    execution_path: list


class Agent:
    def __init__(self, client, model_name, get_user_message, tools: List[ToolDefinition], max_history: int = 15):
        self.client = client
        self.model_name = model_name
        self.get_user_message = get_user_message
        self.tools = tools
        self.max_history = max_history
        self.control = AgentControlState()

        # Convert ToolDefinitions to LangChain tools
        self.langchain_tools = [tool_definition_to_langchain(td) for td in tools]
        self.tool_map = {t.name: t for t in self.langchain_tools}

        # Create ChatOpenAI model and bind tools
        self.llm = ChatOpenAI(
            model=model_name,
            api_key=client.api_key,
        )
        self.llm_with_tools = self.llm.bind_tools(self.langchain_tools)

        # Keep for backward compat
        self.openai_tools = self._convert_tools_to_openai_format()

        # Build the LangGraph
        self._graph = self._build_graph()

        # Test hook: set to a node name to simulate a failure there
        self._test_fail_node = None

        # System instruction for agentic behavior
        self.system_instruction = """
        You are an expert AI software engineer. When performing tasks:
        1. Always verify the state of the filesystem before and after your actions.
        2. If a tool returns an error, analyze the cause and try a different approach.
        3. Use the 'check_syntax', 'run_tests', and 'lint_code' tools to identify and fix errors:
           - Syntax Errors: Use 'check_syntax' to catch compile-time or parse errors.
           - Logical Errors: Use 'run_tests' to verify behavior against expectations.
           - Static Analysis: Use 'lint_code' to find code smells and potential bugs.
        4. Be concise but thorough.
        5. You have FULL ACCESS to the local filesystem using absolute or relative paths. Do not claim you cannot access or retrieve files; instead, use the provided tools (like 'read_file', 'list_files', or specifying paths in tool arguments) to interact with them.
        6. Use absolute or relative paths as provided. Default to the current working directory ('.') if no path is explicitly specified. Do not assume previous paths from history apply to new, unrelated tasks. If you need to switch to a different project or directory, use the 'change_working_directory' tool.
        7. If you cannot find a file or directory in the current directory, use the 'search_file' tool or simply use 'read_file', 'list_files', or 'create_and_edit_file' with the name; the system will automatically search common directories (Desktop, Documents, Pictures, projects, repos, etc.) for you.
        8. If 'read_file' indicates a file is an image, use 'recognize_image' to analyze its contents.
        9. When writing code inside any and all types of code files, write multi-line code and avoid single-line writes.
        10. For SIMPLE tasks (single-step, clear intent, fewer than 3 tool calls): Do not ask for permission; just execute the necessary steps to complete the task. For COMPLEX tasks, follow rules 17-18 instead.
        11. In your final summary, avoid using the words "Error" or "Exception" if the task was completed successfully, as these words are used for automated failure detection. Use words like "issue", "problem", or "fault" if you must refer to them.
        12. ALWAYS report the actual output from tool executions. Never hallucinate or skip reporting the execution results.
        13. If the user denies a tool call or action, explicitly state in your response that you could not finish the task because the user denied it.
        14. For ALL GitHub operations (creating branches, committing files, raising pull requests), use ONLY the github_create_branch, github_commit_file, github_commit_local_file, and github_create_pr MCP tools. NEVER use run_code with git commands for GitHub operations. The workflow is: Step 1: github_create_branch -> Step 2: github_commit_file or github_commit_local_file -> Step 3: github_create_pr.
        15. DOCUMENT PAGE COUNT: When the user requests a document with a specific number of pages or slides, you MUST honor that request exactly:
           - For PDF/DOCX: Set the 'pages' parameter AND insert exactly (N-1) <!-- PAGE_BREAK --> markers in your content for N pages. Write 300-450 words per page to fill each page. Each <!-- PAGE_BREAK --> marker MUST be on its own line.
           - For PPTX: Set the 'pages' parameter AND include EXACTLY N headings (# or ##) for N slides. Count your headings before submitting.
           - NEVER create fewer or more pages/slides than requested. Double-check your content structure before calling the tool.
        16. TRAVEL BOOKING WORKFLOW:
           - For flights: Use 'search_flights' to find available flights. Results include Duffel offer IDs (e.g. [Offer: off_xxx]). When the user selects a flight, collect their details step-by-step per the book_travel tool description, then call 'book_travel' with the offer_id and passenger details.
           - Offers expire in approximately 30 minutes. If booking fails due to expiry, search again.
           - Use 'get_booking' to look up existing bookings by reference, 'cancel_booking' to cancel, and 'list_bookings' to show all bookings.
           - Always confirm the total price and all details with the user before calling book_travel.
        17. TASK CLASSIFICATION — Before acting on any user request, classify it:
           - SIMPLE: The task has clear intent, requires fewer than 3 tool calls, and has no ambiguity in scope or approach (e.g. "read file X", "what is 2+2", "create a hello world script"). For simple tasks, execute immediately per rule 10.
           - COMPLEX: The task is multi-step (3+ tool calls), has ambiguity in scope/approach, touches multiple files, or requires architectural decisions (e.g. "build me a REST API", "refactor the auth system", "add a TUI to this project"). For complex tasks, follow rule 18.
           If unsure, treat it as complex. It is better to ask one unnecessary question than to waste effort building the wrong thing.
        18. COMPLEX TASK PROTOCOL — When a task is classified as COMPLEX, follow these phases IN ORDER. Do NOT call any tools during phases 1 and 2; respond with text only.
           PHASE 1 — CLARIFY: Ask 1-3 focused questions to understand the user's requirements. Ask ONE question per response. Focus on: scope, constraints, preferred approach, and success criteria. Examples: "Should this include authentication?", "Which framework do you prefer?", "What files should I avoid changing?" Move to Phase 2 once you have enough context. SKIP this phase if the user explicitly says "just do it" or provides comprehensive details upfront.
           PHASE 2 — PLAN: Present a numbered step-by-step plan of what you will do. Include which files you will create or modify. End with "Ready to proceed?" and WAIT for the user to confirm before moving to Phase 3. If the user requests changes to the plan, revise and present again.
           PHASE 3 — EXECUTE: Carry out the plan step by step using tools. Follow all other rules (1-16) during execution. After completing each major step, briefly report what was done and what comes next.
           OVERRIDE: If at any point the user says "just do it", "skip the questions", or "go ahead", immediately move to Phase 3 and execute.
        19. PROGRESS REPORTING: When executing a multi-step plan (Phase 3 of rule 18), after each major step report: what you just completed, and what the next step is. Keep progress updates to 1-2 sentences. Example: "Step 2 complete: created models.py with User schema. Next: adding API routes in views.py."
        """

    # ----------------------------------------------------------------
    #  LangGraph: build the execution graph
    # ----------------------------------------------------------------

    def _build_graph(self):
        """Build and compile the LangGraph StateGraph."""
        agent_self = self

        # -- Graph nodes --------------------------------------------------

        def call_model(state: AgentState) -> dict:
            """Invoke the LLM with current messages and bound tools."""
            path = state.get("execution_path", []) + ["call_model"]
            if state["stopped"]:
                return {"status": "stopped", "response": "\u26d4 Execution stopped mid-response.", "execution_path": path}
            try:
                if agent_self._test_fail_node == "call_model":
                    agent_self._test_fail_node = None
                    raise ConnectionError("Simulated failure: Could not reach LLM API (no internet connection)")
                response = agent_self.llm_with_tools.invoke(state["messages"])
                return {"messages": state["messages"] + [response], "execution_path": path}
            except Exception as e:
                error_path = state.get("execution_path", []) + ["call_model \u2717"]
                return {
                    "status": "error",
                    "response": f"LLM call failed at **call_model**: {str(e)}",
                    "execution_path": error_path,
                    "response_history": agent_self._messages_to_dicts(state["messages"]),
                }

        def collect_dry_run(state: AgentState) -> dict:
            """Collect all tool calls into a dry-run plan without executing."""
            path = state.get("execution_path", []) + ["collect_dry_run"]
            try:
                if agent_self._test_fail_node == "collect_dry_run":
                    agent_self._test_fail_node = None
                    raise RuntimeError("Simulated failure: Could not build dry-run plan")
                ai_message = state["messages"][-1]
                plan = []
                for tc in ai_message.tool_calls:
                    plan.append({
                        "id": tc["id"],
                        "name": tc["name"],
                        "arguments": tc["args"],
                        "summary": agent_self._summarize_tool_call(tc["name"], tc["args"])
                    })

                user_request = ""
                for msg in reversed(state["messages"]):
                    if isinstance(msg, HumanMessage):
                        user_request = msg.content
                        break

                history = agent_self._messages_to_dicts(state["messages"])
                return {
                    "status": "dry_run",
                    "dry_run_plan": plan,
                    "response": agent_self._generate_plan_summary(plan, user_request),
                    "response_history": history,
                    "execution_path": path,
                }
            except Exception as e:
                error_path = state.get("execution_path", []) + ["collect_dry_run \u2717"]
                return {
                    "status": "error",
                    "response": f"Failed at **collect_dry_run**: {str(e)}",
                    "execution_path": error_path,
                    "response_history": agent_self._messages_to_dicts(state["messages"]),
                }

        def execute_or_hold_tools(state: AgentState) -> dict:
            """Execute low-risk tools immediately; hold high-risk tools for approval."""
            path = state.get("execution_path", []) + ["execute_or_hold_tools"]
            try:
                if agent_self._test_fail_node == "execute_or_hold_tools":
                    agent_self._test_fail_node = None
                    raise RuntimeError("Simulated failure: Tool execution engine crashed")
                ai_message = state["messages"][-1]
                new_messages = list(state["messages"])
                pending = []

                for tc in ai_message.tool_calls:
                    if agent_self.control.stopped:
                        return {
                            "messages": new_messages,
                            "status": "stopped",
                            "response": "\u26d4 Agent execution was stopped.",
                            "response_history": agent_self._messages_to_dicts(new_messages),
                            "execution_path": path,
                        }

                    tool = agent_self.tool_map.get(tc["name"])
                    if tool and tool.requires_approval:
                        pending.append({
                            "id": tc["id"],
                            "name": tc["name"],
                            "arguments": tc["args"]
                        })
                    else:
                        result = agent_self._execute_tool_by_name(tc["name"], tc["args"])
                        new_messages.append(ToolMessage(
                            content=result,
                            tool_call_id=tc["id"],
                            name=tc["name"],
                        ))

                if pending:
                    history = agent_self._messages_to_dicts(new_messages)
                    return {
                        "messages": new_messages,
                        "status": "pending",
                        "pending_tools": pending,
                        "response": ai_message.content or "Approve the next action:",
                        "response_history": history,
                        "execution_path": path,
                    }

                # All tools were low-risk and already executed; loop back to model
                return {"messages": new_messages, "pending_tools": [], "execution_path": path}
            except Exception as e:
                error_path = state.get("execution_path", []) + ["execute_or_hold_tools \u2717"]
                return {
                    "messages": state["messages"],
                    "status": "error",
                    "response": f"Tool execution failed at **execute_or_hold_tools**: {str(e)}",
                    "execution_path": error_path,
                    "response_history": agent_self._messages_to_dicts(state["messages"]),
                }

        def format_output(state: AgentState) -> dict:
            """Format the final text response."""
            path = state.get("execution_path", []) + ["format_output"]
            try:
                if agent_self._test_fail_node == "format_output":
                    agent_self._test_fail_node = None
                    raise RuntimeError("Simulated failure: Could not format output")
                history = agent_self._messages_to_dicts(state["messages"])
                result = {"response_history": history, "execution_path": path}
                # If status was already set (e.g. "stopped" or "error"), preserve it
                if not state.get("status"):
                    last = state["messages"][-1]
                    result["status"] = "success"
                    result["response"] = (
                        last.content if isinstance(last, AIMessage) else "No response generated"
                    ) or "No response generated"
                return result
            except Exception as e:
                error_path = state.get("execution_path", []) + ["format_output \u2717"]
                return {
                    "status": "error",
                    "response": f"Failed at **format_output**: {str(e)}",
                    "execution_path": error_path,
                    "response_history": [],
                }

        # -- Conditional routing ------------------------------------------

        def route_after_model(state: AgentState) -> str:
            if state.get("status") in ("stopped", "error"):
                return "format_output"
            last = state["messages"][-1]
            if not isinstance(last, AIMessage) or not last.tool_calls:
                return "format_output"
            if state["use_pending"]:
                return "execute_or_hold_tools"
            return "collect_dry_run"

        def route_after_tools(state: AgentState) -> str:
            if state.get("status") in ("stopped", "error"):
                return "format_output"
            if state["pending_tools"]:
                return END
            return "call_model"

        # -- Assemble graph -----------------------------------------------

        graph = StateGraph(AgentState)
        graph.add_node("call_model", call_model)
        graph.add_node("collect_dry_run", collect_dry_run)
        graph.add_node("execute_or_hold_tools", execute_or_hold_tools)
        graph.add_node("format_output", format_output)

        graph.set_entry_point("call_model")
        graph.add_conditional_edges("call_model", route_after_model, {
            "format_output": "format_output",
            "collect_dry_run": "collect_dry_run",
            "execute_or_hold_tools": "execute_or_hold_tools",
        })
        graph.add_conditional_edges("execute_or_hold_tools", route_after_tools, {
            END: END,
            "call_model": "call_model",
            "format_output": "format_output",
        })
        graph.add_edge("collect_dry_run", END)
        graph.add_edge("format_output", END)

        return graph.compile()

    # ----------------------------------------------------------------
    #  Public API (same signatures as the original Agent)
    # ----------------------------------------------------------------

    def chat_once(self, conversation_history=None, message=None, use_pending=False):
        """
        Handle a single chat interaction for Django/API usage.

        Args:
            conversation_history: List of previous messages (optional)
            message: Single message string to process
            use_pending: If True the model's tool calls go through per-tool
                         approval (status="pending").  If False (the default,
                         used for every fresh user message) all tool calls are
                         collected into a single dry-run plan first.

        Returns:
            Dict containing status and response/tool info
        """
        try:
            # Reset stopped flag for new messages to allow continuation
            if message is not None:
                self.control.stopped = False

            if self.control.stopped:
                return {
                    "status": "stopped",
                    "response": "\u26d4 Agent execution was stopped."
                }

            # Build LangChain message list
            messages = [SystemMessage(content=self.system_instruction)]
            if conversation_history:
                for lc_msg in self._dicts_to_messages(conversation_history):
                    if isinstance(lc_msg, SystemMessage):
                        continue
                    messages.append(lc_msg)

            if message:
                if is_prompt_injection(message):
                    return {"status": "error", "message": "Security Warning: Potential prompt injection detected. Message blocked."}
                messages.append(HumanMessage(content=message))

            messages = self._trim_messages(messages)

            initial_state: AgentState = {
                "messages": messages,
                "use_pending": use_pending,
                "dry_run_plan": [],
                "pending_tools": [],
                "status": "",
                "response": "",
                "response_history": [],
                "stopped": self.control.stopped,
                "tools_enabled": self.control.tools_enabled,
                "execution_path": ["__start__"],
            }

            result = self._graph.invoke(initial_state, {"recursion_limit": 25})

            # Build execution path string
            path = result.get("execution_path", []) + ["__end__"]

            # Build output in the same format the frontend expects
            output = {
                "status": result["status"],
                "response": result["response"],
                "execution_path": path,
            }
            if result["status"] == "dry_run":
                output["dry_run_plan"] = result["dry_run_plan"]
                output["history"] = result["response_history"]
            elif result["status"] == "pending":
                output["pending_tools"] = result["pending_tools"]
                output["history"] = result["response_history"]
            else:
                output["history"] = result["response_history"]

            return output

        except Exception as e:
            return {
                "status": "error",
                "response": f"Error: {str(e)}",
                "message": str(e),
                "execution_path": ["__start__", "graph_error \u2717", "__end__"],
            }

    def execute_dry_run(self, dry_run_plan: List[Dict[str, Any]], history: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Execute every tool in an approved dry-run plan, then fetch the model follow-up.

        The follow-up is processed with use_pending=True so any further tool
        calls the model issues go through per-tool approval rather than
        surfacing as another dry-run round.
        """
        try:
            # Convert dict history to LangChain messages (history already contains system msg)
            messages = self._dicts_to_messages(history)
            if not messages or not isinstance(messages[0], SystemMessage):
                messages.insert(0, SystemMessage(content=self.system_instruction))

            # Execute each approved tool
            for tool_call in dry_run_plan:
                if self.control.stopped:
                    return {"status": "stopped", "response": "\u26d4 Agent execution was stopped."}

                result = self._execute_tool_by_name(tool_call['name'], tool_call['arguments'])
                messages.append(ToolMessage(
                    content=result,
                    tool_call_id=tool_call['id'],
                    name=tool_call['name'],
                ))

            # Run the graph for follow-up in per-tool mode
            initial_state: AgentState = {
                "messages": self._trim_messages(messages),
                "use_pending": True,
                "dry_run_plan": [],
                "pending_tools": [],
                "status": "",
                "response": "",
                "response_history": [],
                "stopped": self.control.stopped,
                "tools_enabled": self.control.tools_enabled,
                "execution_path": ["__start__"],
            }

            result = self._graph.invoke(initial_state, {"recursion_limit": 25})

            path = result.get("execution_path", []) + ["__end__"]

            output = {
                "status": result["status"],
                "response": result["response"],
                "execution_path": path,
            }
            if result["status"] == "pending":
                output["pending_tools"] = result["pending_tools"]
            output["history"] = result["response_history"]
            return output

        except Exception as e:
            return {
                "status": "error",
                "response": f"Error executing dry run: {str(e)}",
                "message": f"Error executing dry run: {str(e)}",
                "execution_path": ["__start__", "execute_dry_run \u2717", "__end__"],
            }

    # ----------------------------------------------------------------
    #  CLI interactive mode
    # ----------------------------------------------------------------

    def run(self):
        print("Chat with OpenAI (use 'ctrl-c' or type 'quit' to exit)")

        self.messages = [SystemMessage(content=self.system_instruction)]

        while True:
            if self.control.stopped:
                print("\u26d4 Agent execution stopped.")
                break

            print("\033[94mYou\033[0m: ", end="")
            user_input, ok = self.get_user_message()
            if not ok:
                break

            # Check for quit command
            if user_input.lower().strip() in ['quit', 'exit', 'q']:
                print("Goodbye!")
                break

            self.messages.append(HumanMessage(content=user_input))

            try:
                response = self.llm_with_tools.invoke(
                    self._trim_messages(self.messages)
                )
                self._process_response_simple(response)
            except Exception as e:
                print(f"Error: {str(e)}")
                import traceback
                print(f"Traceback: {traceback.format_exc()}")

    def _process_response_simple(self, ai_message):
        """Simplified response processing for CLI mode (no approval system)."""
        try:
            self.messages.append(ai_message)

            # Handle text response
            if ai_message.content:
                print(f"\033[93mOpenAI\033[0m: {ai_message.content}")

            # Handle tool calls
            if ai_message.tool_calls:
                for tc in ai_message.tool_calls:
                    args_str = json.dumps(tc["args"]).replace('\\\\n', '\\n').replace('\\n', '\n')
                    print(f"\033[92mtool\033[0m: {tc['name']}({args_str})")

                    result = self._execute_tool_by_name(tc["name"], tc["args"])

                    self.messages.append(ToolMessage(
                        content=result,
                        tool_call_id=tc["id"],
                        name=tc["name"],
                    ))

                # Get follow-up response
                try:
                    follow_up = self.llm_with_tools.invoke(
                        self._trim_messages(self.messages)
                    )
                    self._process_response_simple(follow_up)
                except Exception as e:
                    print(f"Error processing follow-up: {str(e)}")

        except Exception as e:
            print(f"Error processing response: {str(e)}")
            import traceback
            print(f"Traceback: {traceback.format_exc()}")

    # ----------------------------------------------------------------
    #  Tool execution (unchanged logic - calls original ToolDefinition.function)
    # ----------------------------------------------------------------

    def _execute_tool_by_name(self, name, args):
        """Find and execute a tool by name."""
        if self.control.stopped:
            return "\u26d4 Agent execution stopped."

        if not self.control.tools_enabled:
            return "\U0001f512 Tool execution disabled by kill switch."

        from chat.models import ToolLog
        args_str = json.dumps(args).replace('\\\\n', '\\n').replace('\\n', '\n')
        print(f"\033[96mTool Call:\033[0m {name}({args_str})")
        for tool in self.tools:
            if tool.name == name:
                try:
                    result = tool.function(args)
                    # Save log to database
                    ToolLog.objects.create(
                        tool_name=name,
                        input_args=json.dumps(args),
                        output_result=str(result)
                    )
                    return result
                except Exception as e:
                    error_msg = f"Error executing tool: {str(e)}"
                    # Save error log to database
                    ToolLog.objects.create(
                        tool_name=name,
                        input_args=json.dumps(args),
                        output_result=error_msg
                    )
                    return error_msg
        return f"Tool '{name}' not found"

    # ----------------------------------------------------------------
    #  Message format conversion helpers
    # ----------------------------------------------------------------

    def _dicts_to_messages(self, history: List[Dict[str, Any]]) -> List[BaseMessage]:
        """Convert frontend dict-format history to LangChain message objects."""
        messages = []
        for msg in history:
            role = msg.get("role")
            if role == "system":
                messages.append(SystemMessage(content=msg.get("content", "")))
            elif role == "user":
                messages.append(HumanMessage(content=msg.get("content", "")))
            elif role == "assistant":
                kwargs: Dict[str, Any] = {"content": msg.get("content") or ""}
                if msg.get("tool_calls"):
                    kwargs["tool_calls"] = [
                        {
                            "id": tc["id"],
                            "name": tc["function"]["name"],
                            "args": json.loads(tc["function"]["arguments"])
                                    if isinstance(tc["function"]["arguments"], str)
                                    else tc["function"]["arguments"],
                        }
                        for tc in msg["tool_calls"]
                    ]
                messages.append(AIMessage(**kwargs))
            elif role == "tool":
                messages.append(ToolMessage(
                    content=msg.get("content", ""),
                    tool_call_id=msg.get("tool_call_id", ""),
                    name=msg.get("name", ""),
                ))
        return messages

    def _messages_to_dicts(self, messages: List[BaseMessage]) -> List[Dict[str, Any]]:
        """Convert LangChain messages back to the frontend dict format."""
        result = []
        for msg in messages:
            if isinstance(msg, SystemMessage):
                result.append({"role": "system", "content": msg.content})
            elif isinstance(msg, HumanMessage):
                result.append({"role": "user", "content": msg.content})
            elif isinstance(msg, AIMessage):
                d: Dict[str, Any] = {"role": "assistant", "content": msg.content}
                if msg.tool_calls:
                    d["tool_calls"] = [
                        {
                            "id": tc["id"],
                            "type": "function",
                            "function": {
                                "name": tc["name"],
                                "arguments": json.dumps(tc["args"])
                            }
                        }
                        for tc in msg.tool_calls
                    ]
                result.append(d)
            elif isinstance(msg, ToolMessage):
                result.append({
                    "role": "tool",
                    "tool_call_id": msg.tool_call_id,
                    "name": msg.name,
                    "content": msg.content,
                })
        return result

    # ----------------------------------------------------------------
    #  Dry-run plan helpers (unchanged logic)
    # ----------------------------------------------------------------

    def _summarize_tool_call(self, name: str, args: Dict[str, Any]) -> str:
        """Return a one-line, human-readable description of a tool call for the dry-run plan."""
        summaries = {
            'read_file':                lambda a: f"Read file: {a.get('path', '')}",
            'list_files':               lambda a: f"List files in: {a.get('path', '.')}",
            'search_file':              lambda a: f"Search for file: {a.get('filename', '')}",
            'find_file_broadly':        lambda a: f"Find file: {a.get('filename', '')}",
            'find_directory_broadly':   lambda a: f"Find directory: {a.get('dirname', '')}",
            'create_and_edit_file':     lambda a: f"Edit/create file: {a.get('path', '')}",
            'delete_file':              lambda a: f"Delete: {a.get('path', '')}",
            'rename_file':              lambda a: f"Rename {a.get('old_path', '')} \u2192 {a.get('new_path', '')}",
            'run_code':                 lambda a: f"Run command: {a.get('command', '')}",
            'check_syntax':             lambda a: f"Check syntax of: {a.get('path', '')}",
            'run_tests':                lambda a: f"Run tests: {a.get('command', 'auto-detect')}",
            'lint_code':                lambda a: f"Lint: {a.get('path', '')}",
            'change_working_directory': lambda a: f"Change directory to: {a.get('path', '')}",
            'create_pdf':               lambda a: f"Create PDF: {a.get('filename', '')}",
            'create_docx':              lambda a: f"Create DOCX: {a.get('filename', '')}",
            'create_excel':             lambda a: f"Create Excel: {a.get('filename', '')}",
            'create_pptx':              lambda a: f"Create PPTX: {a.get('filename', '')}",
            'open_gmail_and_compose':   lambda a: f"Compose email to: {a.get('recipient', '')}",
            'recognize_image':          lambda a: f"Analyze image: {a.get('path', '')}",
            'recognize_video':          lambda a: f"Analyze video: {a.get('path', '')}",
            'recognize_audio':          lambda a: f"Analyze audio: {a.get('path', '')}",
            'github_create_branch':     lambda a: f"Create GitHub branch: {a.get('name', '')}",
            'github_commit_file':       lambda a: f"Commit to GitHub branch '{a.get('branch', '')}': {a.get('path', '')}",
            'github_commit_local_file': lambda a: f"Commit local file to GitHub branch '{a.get('branch', '')}': {a.get('local_path', '')}",
            'github_create_pr':         lambda a: f'Create GitHub PR: "{a.get("title", "")}"',
            'playwright_navigate':      lambda a: f"Navigate to: {a.get('url', '')}",
        }
        summarizer = summaries.get(name)
        return summarizer(args) if summarizer else f"Call '{name}' with: {json.dumps(args)}"

    def _generate_plan_summary(self, dry_run_plan: List[Dict[str, Any]], user_request: str = "") -> str:
        """Ask the model to produce a plain-English, start-to-finish summary
        of the dry-run plan.  The user's original request is included so the
        summary stays grounded.  Falls back to a generic string on any failure
        so the card still renders."""
        steps = "\n".join(
            f"{i + 1}. {step['summary']}  \u2014  {step['name']}({json.dumps(step['arguments'])})"
            for i, step in enumerate(dry_run_plan)
        )
        user_ctx = f"The user asked: \"{user_request}\"\n\n" if user_request else ""

        try:
            resp = self.llm.invoke([
                SystemMessage(content=(
                    "You are about to execute an action plan on behalf of the user. "
                    "Summarise the plan below in 2-4 clear sentences of plain English. "
                    "Describe what will happen from start to finish and what the end "
                    "result will be. Do not repeat raw tool names or JSON \u2014 translate "
                    "everything into natural language."
                )),
                HumanMessage(content=f"{user_ctx}Planned actions:\n\n{steps}\n\nSummarise what this plan will do:")
            ])
            summary = resp.content
            if summary and summary.strip():
                return summary.strip()
        except Exception as e:
            print(f"[dry_run] summary generation failed: {e}")

        return "Here is the planned action(s). Review and approve or deny before execution:"

    # ----------------------------------------------------------------
    #  Backward-compat helpers
    # ----------------------------------------------------------------

    def _convert_tools_to_openai_format(self):
        """Convert tools to OpenAI format (kept for backward compatibility)."""
        openai_tools = []
        for tool in self.tools:
            openai_tools.append({
                "type": "function",
                "function": {
                    "name": tool.name,
                    "description": tool.description,
                    "parameters": tool.parameters
                }
            })
        return openai_tools

    def _trim_messages(self, messages: List[Any]) -> List[Any]:
        """
        Trim message history to stay within token limits.
        Keeps the system message and ensures tool messages are not separated from their tool_calls.
        Works with both LangChain message objects and plain dicts.
        """
        if len(messages) <= self.max_history + 1:
            return messages

        first = messages[0]
        is_system = (
            isinstance(first, SystemMessage)
            or (isinstance(first, dict) and first.get("role") == "system")
        )
        system_message = first if is_system else None

        # Keep the most recent messages
        start_index = len(messages) - self.max_history

        # Ensure we don't start in the middle of a tool response sequence
        def _is_tool(m):
            if isinstance(m, ToolMessage):
                return True
            if isinstance(m, dict) and m.get("role") == "tool":
                return True
            return False

        while start_index > 1 and _is_tool(messages[start_index]):
            start_index -= 1

        recent_messages = messages[start_index:]

        trimmed = []
        if system_message:
            trimmed.append(system_message)
        trimmed.extend(recent_messages)

        return trimmed

    def generate_code(self, task_description: str, language: str = "python", stepwise: bool = True) -> str:
        """
        Generate complex code using the LLM with advanced prompt engineering.
        Args:
            task_description (str): Description of the code to generate.
            language (str): Programming language for the code.
            stepwise (bool): Whether to use a stepwise prompt for better results.
        Returns:
            str: Generated code.
        """
        if is_prompt_injection(task_description):
            return "Security Warning: Potential prompt injection detected in task description. Generation blocked."

        if stepwise:
            prompt = (
                f"""
                Write a complete, well-structured {language} program for the following task:
                {task_description}

                Please follow these steps:
                1. Start by outlining the main components or functions needed.
                2. Implement each component step by step, with clear comments.
                3. At the end, provide the full code in a single code block.
                4. Ensure the code is ready to run and includes all necessary parts (imports, main function, etc.).
                """
            )
        else:
            prompt = f"Write a complete {language} program for the following task: {task_description}"

        response = self.llm.invoke([HumanMessage(content=prompt)])
        return response.content.strip()


if __name__ == "__main__":
    main()
