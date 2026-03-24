import os
import re
from typing import Dict, Any
from agents.helpers import find_file_broadly
from agents.control import ToolDefinition
from agents.document_tools import create_pdf, create_docx, create_excel, create_pptx

from docx import Document
from openpyxl import Workbook
from pptx import Presentation


# ─── Document Read/Edit Functions ────────────────────────────────────

def read_pdf_tool(args: Dict[str, Any]) -> str:
    """Read the text content of a PDF file."""
    path = args.get('path', '')
    if not path:
        return "Error: No path provided."

    actual_path = path
    if not os.path.exists(path):
        found = find_file_broadly(path)
        if found:
            actual_path = found
        else:
            return f"Error: File '{path}' not found."

    try:
        import fitz  # PyMuPDF
        doc = fitz.open(actual_path)
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                pages.append(f"--- Page {i + 1} ---\n{text.strip()}")
        doc.close()

        if not pages:
            return f"The PDF '{actual_path}' contains no extractable text (may be image-based)."

        result = f"Content of {actual_path} ({len(pages)} page(s)):\n\n" + "\n\n".join(pages)
        if len(result) > 15000:
            result = result[:15000] + "\n\n[... Output truncated. Use a page range for large PDFs ...]"
        return result
    except ImportError:
        # Fallback to pdfplumber if PyMuPDF not available
        try:
            import pdfplumber
            pages = []
            with pdfplumber.open(actual_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    if text and text.strip():
                        pages.append(f"--- Page {i + 1} ---\n{text.strip()}")
            if not pages:
                return f"The PDF '{actual_path}' contains no extractable text."
            result = f"Content of {actual_path} ({len(pages)} page(s)):\n\n" + "\n\n".join(pages)
            if len(result) > 15000:
                result = result[:15000] + "\n\n[... Output truncated ...]"
            return result
        except ImportError:
            return "Error: Neither PyMuPDF (fitz) nor pdfplumber is installed. Install one with: pip install PyMuPDF or pip install pdfplumber"
    except Exception as e:
        return f"Error reading PDF: {str(e)}"


def read_docx_tool(args: Dict[str, Any]) -> str:
    """Read the text content of a DOCX file including paragraphs and tables."""
    path = args.get('path', '')
    if not path:
        return "Error: No path provided."

    actual_path = path
    if not os.path.exists(path):
        found = find_file_broadly(path)
        if found:
            actual_path = found
        else:
            return f"Error: File '{path}' not found."

    try:
        doc = Document(actual_path)
        parts = []

        for element in doc.element.body:
            tag = element.tag.split('}')[-1] if '}' in element.tag else element.tag

            if tag == 'p':
                # Paragraph
                para = None
                for p in doc.paragraphs:
                    if p._element is element:
                        para = p
                        break
                if para and para.text.strip():
                    style_name = para.style.name if para.style else ''
                    if 'Heading' in style_name:
                        level = ''.join(c for c in style_name if c.isdigit()) or '1'
                        parts.append(f"{'#' * int(level)} {para.text}")
                    elif 'List' in style_name:
                        parts.append(f"- {para.text}")
                    else:
                        parts.append(para.text)

            elif tag == 'tbl':
                # Table
                for table in doc.tables:
                    if table._element is element:
                        rows = []
                        for ri, row in enumerate(table.rows):
                            cells = [cell.text.strip() for cell in row.cells]
                            rows.append("| " + " | ".join(cells) + " |")
                            if ri == 0:
                                rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                        parts.append("\n".join(rows))
                        break

        if not parts:
            return f"The DOCX '{actual_path}' appears to be empty."

        result = f"Content of {actual_path}:\n\n" + "\n\n".join(parts)
        if len(result) > 15000:
            result = result[:15000] + "\n\n[... Output truncated ...]"
        return result
    except Exception as e:
        return f"Error reading DOCX: {str(e)}"


def read_excel_tool(args: Dict[str, Any]) -> str:
    """Read the content of an XLSX file, returning data as markdown tables."""
    path = args.get('path', '')
    sheet_name = args.get('sheet_name', '')
    if not path:
        return "Error: No path provided."

    actual_path = path
    if not os.path.exists(path):
        found = find_file_broadly(path)
        if found:
            actual_path = found
        else:
            return f"Error: File '{path}' not found."

    try:
        wb = Workbook()
        wb.close()
        from openpyxl import load_workbook
        wb = load_workbook(actual_path, read_only=True, data_only=True)

        sheets_to_read = []
        if sheet_name:
            if sheet_name in wb.sheetnames:
                sheets_to_read = [sheet_name]
            else:
                wb.close()
                return f"Error: Sheet '{sheet_name}' not found. Available sheets: {', '.join(wb.sheetnames)}"
        else:
            sheets_to_read = wb.sheetnames

        parts = []
        for sn in sheets_to_read:
            ws = wb[sn]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_vals = [str(c) if c is not None else "" for c in row]
                if any(v.strip() for v in row_vals):
                    rows.append(row_vals)

            if not rows:
                parts.append(f"### Sheet: {sn}\n(empty)")
                continue

            # Format as markdown table
            col_count = max(len(r) for r in rows)
            # Pad rows to same length
            for r in rows:
                while len(r) < col_count:
                    r.append("")

            md_lines = []
            md_lines.append("| " + " | ".join(rows[0]) + " |")
            md_lines.append("| " + " | ".join(["---"] * col_count) + " |")
            for r in rows[1:]:
                md_lines.append("| " + " | ".join(r) + " |")

            if len(sheets_to_read) > 1:
                parts.append(f"### Sheet: {sn}\n" + "\n".join(md_lines))
            else:
                parts.append("\n".join(md_lines))

        wb.close()

        result = f"Content of {actual_path} ({len(sheets_to_read)} sheet(s)):\n\n" + "\n\n".join(parts)
        if len(result) > 15000:
            result = result[:15000] + "\n\n[... Output truncated. Specify a sheet_name to read a specific sheet ...]"
        return result
    except Exception as e:
        return f"Error reading XLSX: {str(e)}"


def read_pptx_tool(args: Dict[str, Any]) -> str:
    """Read the text content of a PPTX file, returning slides as markdown."""
    path = args.get('path', '')
    if not path:
        return "Error: No path provided."

    actual_path = path
    if not os.path.exists(path):
        found = find_file_broadly(path)
        if found:
            actual_path = found
        else:
            return f"Error: File '{path}' not found."

    try:
        prs = Presentation(actual_path)
        parts = []

        for si, slide in enumerate(prs.slides, 1):
            slide_parts = [f"## Slide {si}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            level = para.level if hasattr(para, 'level') else 0
                            if level > 0:
                                slide_parts.append(f"{'  ' * level}- {text}")
                            else:
                                slide_parts.append(text)
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for ri, row in enumerate(table.rows):
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append("| " + " | ".join(cells) + " |")
                        if ri == 0:
                            rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                    slide_parts.append("\n".join(rows))

            parts.append("\n".join(slide_parts))

        if not parts:
            return f"The PPTX '{actual_path}' appears to be empty."

        result = f"Content of {actual_path} ({len(prs.slides)} slide(s)):\n\n" + "\n\n".join(parts)
        if len(result) > 15000:
            result = result[:15000] + "\n\n[... Output truncated ...]"
        return result
    except Exception as e:
        return f"Error reading PPTX: {str(e)}"


def edit_pdf_tool(args: Dict[str, Any]) -> str:
    """Edit a PDF by reading its content, applying changes, and rewriting it."""
    path = args.get('path', '')
    content = args.get('content', '')

    if not path:
        return "Error: No path provided."
    if not content.strip():
        return "Error: No content provided. Provide the full updated document content in markdown format."

    actual_path = path
    if not os.path.exists(path):
        found = find_file_broadly(path)
        if found:
            actual_path = found
        else:
            return f"Error: File '{path}' not found."

    return create_pdf(actual_path, content)


def edit_docx_tool(args: Dict[str, Any]) -> str:
    """Edit a DOCX by reading its content, applying changes, and rewriting it."""
    path = args.get('path', '')
    content = args.get('content', '')

    if not path:
        return "Error: No path provided."
    if not content.strip():
        return "Error: No content provided. Provide the full updated document content in markdown format."

    actual_path = path
    if not os.path.exists(path):
        found = find_file_broadly(path)
        if found:
            actual_path = found
        else:
            return f"Error: File '{path}' not found."

    return create_docx(actual_path, content)


def edit_excel_tool(args: Dict[str, Any]) -> str:
    """Edit an XLSX file by rewriting it with new data."""
    path = args.get('path', '')
    content = args.get('content', '')
    data = args.get('data')
    title = args.get('title', '')

    if not path:
        return "Error: No path provided."
    if not content.strip() and not data:
        return "Error: No content or data provided."

    actual_path = path
    if not os.path.exists(path):
        found = find_file_broadly(path)
        if found:
            actual_path = found
        else:
            return f"Error: File '{path}' not found."

    return create_excel(actual_path, content, data, title)


def edit_pptx_tool(args: Dict[str, Any]) -> str:
    """Edit a PPTX by reading its content, applying changes, and rewriting it."""
    path = args.get('path', '')
    content = args.get('content', '')

    if not path:
        return "Error: No path provided."
    if not content.strip():
        return "Error: No content provided. Provide the full updated presentation content in markdown format."

    actual_path = path
    if not os.path.exists(path):
        found = find_file_broadly(path)
        if found:
            actual_path = found
        else:
            return f"Error: File '{path}' not found."

    return create_pptx(actual_path, content)


# ─── Document Read/Edit Tool Definitions ─────────────────────────────

READ_PDF_DEFINITION = ToolDefinition(
    name="read_pdf",
    description=(
        "Read and extract the text content of a PDF file. Returns all pages as formatted text. "
        "Use this to read existing PDF documents before editing or summarizing them."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the PDF file to read."
            }
        },
        "required": ["path"]
    },
    function=read_pdf_tool,
    requires_approval=False
)

READ_DOCX_DEFINITION = ToolDefinition(
    name="read_docx",
    description=(
        "Read and extract the text content of a Word document (DOCX). Returns headings, paragraphs, "
        "lists, and tables as formatted markdown text. Use this to read existing DOCX files before editing."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the DOCX file to read."
            }
        },
        "required": ["path"]
    },
    function=read_docx_tool,
    requires_approval=False
)

READ_EXCEL_DEFINITION = ToolDefinition(
    name="read_excel",
    description=(
        "Read and extract the data from an Excel spreadsheet (XLSX). Returns each sheet's data as a "
        "markdown table. Optionally specify a sheet_name to read only that sheet."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the XLSX file to read."
            },
            "sheet_name": {
                "type": "string",
                "description": "Optional: name of a specific sheet to read. If omitted, reads all sheets."
            }
        },
        "required": ["path"]
    },
    function=read_excel_tool,
    requires_approval=False
)

READ_PPTX_DEFINITION = ToolDefinition(
    name="read_pptx",
    description=(
        "Read and extract the text content of a PowerPoint presentation (PPTX). Returns each slide's "
        "text content including titles, bullet points, and tables as markdown."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the PPTX file to read."
            }
        },
        "required": ["path"]
    },
    function=read_pptx_tool,
    requires_approval=False
)

EDIT_PDF_DEFINITION = ToolDefinition(
    name="edit_pdf",
    description=(
        "Edit an existing PDF file by rewriting it with updated content. First use read_pdf to get the "
        "current content, then modify it and pass the full updated content in markdown format. "
        "The entire document is regenerated with professional styling."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the PDF file to edit."
            },
            "content": {
                "type": "string",
                "description": "The FULL updated document content in markdown format."
            }
        },
        "required": ["path", "content"]
    },
    function=edit_pdf_tool,
    requires_approval=True
)

EDIT_DOCX_DEFINITION = ToolDefinition(
    name="edit_docx",
    description=(
        "Edit an existing Word document (DOCX) by rewriting it with updated content. First use read_docx "
        "to get the current content, then modify it and pass the full updated content in markdown format. "
        "The entire document is regenerated with professional styling."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the DOCX file to edit."
            },
            "content": {
                "type": "string",
                "description": "The FULL updated document content in markdown format."
            }
        },
        "required": ["path", "content"]
    },
    function=edit_docx_tool,
    requires_approval=True
)

EDIT_EXCEL_DEFINITION = ToolDefinition(
    name="edit_excel",
    description=(
        "Edit an existing Excel spreadsheet (XLSX) by rewriting it with updated data. First use read_excel "
        "to get the current content, then pass updated data as a markdown table in 'content' or as an "
        "array of arrays in 'data'."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the XLSX file to edit."
            },
            "content": {
                "type": "string",
                "description": "Updated data as a markdown table."
            },
            "data": {
                "type": "array",
                "items": {"type": "array", "items": {"type": "string"}},
                "description": "Updated data as array of arrays. First row is headers."
            },
            "title": {
                "type": "string",
                "description": "Optional title row above the data."
            }
        },
        "required": ["path"]
    },
    function=edit_excel_tool,
    requires_approval=True
)

EDIT_PPTX_DEFINITION = ToolDefinition(
    name="edit_pptx",
    description=(
        "Edit an existing PowerPoint presentation (PPTX) by rewriting it with updated content. First use "
        "read_pptx to get the current content, then modify it and pass the full updated content in markdown "
        "format. Each # or ## heading starts a new slide."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "The path to the PPTX file to edit."
            },
            "content": {
                "type": "string",
                "description": "The FULL updated presentation content in markdown format. Each # or ## heading starts a new slide."
            }
        },
        "required": ["path", "content"]
    },
    function=edit_pptx_tool,
    requires_approval=True
)
