import os
import re
from typing import Dict, Any
from agents.helpers import find_file_broadly
from agents.control import ToolDefinition

from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, HRFlowable
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors as rl_colors
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_LEFT, TA_CENTER
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from openpyxl import Workbook
from openpyxl.styles import Font as XlFont, PatternFill, Alignment as XlAlignment, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN


def _parse_markdown_blocks(text):
    """Parse markdown text into structured blocks for document rendering."""
    blocks = []
    lines = text.split('\n')
    i = 0

    while i < len(lines):
        line = lines[i]

        # Page break marker
        if line.strip() in ('<!-- PAGE_BREAK -->', '<!--PAGE_BREAK-->', '---PAGE_BREAK---'):
            blocks.append({'type': 'page_break'})
            i += 1
            continue

        # Code block (fenced)
        if line.strip().startswith('```'):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('```'):
                code_lines.append(lines[i])
                i += 1
            blocks.append({'type': 'code', 'text': '\n'.join(code_lines)})
            i += 1
            continue

        # Heading
        heading_m = re.match(r'^(#{1,4})\s+(.+)$', line)
        if heading_m:
            blocks.append({'type': 'heading', 'level': len(heading_m.group(1)), 'text': heading_m.group(2).strip()})
            i += 1
            continue

        # Horizontal rule
        if re.match(r'^[-*_]{3,}\s*$', line.strip()):
            blocks.append({'type': 'hr'})
            i += 1
            continue

        # Table
        if '|' in line and i + 1 < len(lines) and re.match(r'^[\s|:-]+$', lines[i + 1]):
            headers = [c.strip() for c in line.strip().strip('|').split('|')]
            i += 2  # skip header and separator
            rows = []
            while i < len(lines) and '|' in lines[i] and lines[i].strip():
                row = [c.strip() for c in lines[i].strip().strip('|').split('|')]
                rows.append(row)
                i += 1
            blocks.append({'type': 'table', 'headers': headers, 'rows': rows})
            continue

        # Bullet list
        if re.match(r'^[\s]*[-*+]\s+', line):
            items = []
            while i < len(lines) and re.match(r'^[\s]*[-*+]\s+', lines[i]):
                item_text = re.sub(r'^[\s]*[-*+]\s+', '', lines[i])
                indent = len(lines[i]) - len(lines[i].lstrip())
                level = min(indent // 2, 2)
                items.append({'text': item_text, 'level': level})
                i += 1
            blocks.append({'type': 'bullet_list', 'items': items})
            continue

        # Numbered list
        if re.match(r'^[\s]*\d+[.)]\s+', line):
            items = []
            while i < len(lines) and re.match(r'^[\s]*\d+[.)]\s+', lines[i]):
                item_text = re.sub(r'^[\s]*\d+[.)]\s+', '', lines[i])
                indent = len(lines[i]) - len(lines[i].lstrip())
                level = min(indent // 2, 2)
                items.append({'text': item_text, 'level': level})
                i += 1
            blocks.append({'type': 'numbered_list', 'items': items})
            continue

        # Paragraph (non-empty line)
        if line.strip():
            para_lines = []
            while i < len(lines) and lines[i].strip() and not re.match(r'^#{1,4}\s+', lines[i]) and not re.match(r'^[-*_]{3,}\s*$', lines[i].strip()) and not re.match(r'^[\s]*[-*+]\s+', lines[i]) and not re.match(r'^[\s]*\d+[.)]\s+', lines[i]) and not lines[i].strip().startswith('```') and not ('|' in lines[i] and i + 1 < len(lines) and re.match(r'^[\s|:-]+$', lines[i + 1])):
                para_lines.append(lines[i])
                i += 1
            blocks.append({'type': 'paragraph', 'text': ' '.join(para_lines)})
            continue

        i += 1

    return blocks


def _md_inline_to_html(text):
    """Convert markdown inline formatting to HTML for reportlab Paragraph."""
    text = re.sub(r'\*\*\*(.+?)\*\*\*', r'<b><i>\1</i></b>', text)
    text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)
    text = re.sub(r'\*(.+?)\*', r'<i>\1</i>', text)
    text = re.sub(r'`(.+?)`', r'<font face="Courier" size="9">\1</font>', text)
    return text


# ─── Document Creation Functions ─────────────────────────────────────

def create_pdf(filename: str, content: str = "", pages: int = None) -> str:
    """Create a professionally styled PDF from markdown content."""
    try:
        if not content.strip():
            return "Error: content is required to create a PDF."

        doc = SimpleDocTemplate(
            filename, pagesize=letter,
            leftMargin=0.75 * inch, rightMargin=0.75 * inch,
            topMargin=0.75 * inch, bottomMargin=0.75 * inch
        )

        styles = getSampleStyleSheet()
        # Custom styles
        styles.add(ParagraphStyle(
            'DocTitle', parent=styles['Title'],
            fontSize=24, textColor=rl_colors.HexColor('#1a237e'),
            spaceAfter=20, alignment=TA_CENTER
        ))
        styles.add(ParagraphStyle(
            'H1', parent=styles['Heading1'],
            fontSize=18, textColor=rl_colors.HexColor('#1a237e'),
            spaceBefore=18, spaceAfter=10, borderWidth=1,
            borderColor=rl_colors.HexColor('#1a237e'), borderPadding=4
        ))
        styles.add(ParagraphStyle(
            'H2', parent=styles['Heading2'],
            fontSize=15, textColor=rl_colors.HexColor('#283593'),
            spaceBefore=14, spaceAfter=8
        ))
        styles.add(ParagraphStyle(
            'H3', parent=styles['Heading3'],
            fontSize=12, textColor=rl_colors.HexColor('#3949ab'),
            spaceBefore=10, spaceAfter=6
        ))
        styles.add(ParagraphStyle(
            'BodyCustom', parent=styles['Normal'],
            fontSize=10.5, leading=15, spaceAfter=8
        ))
        styles.add(ParagraphStyle(
            'CodeBlock', parent=styles['Normal'],
            fontName='Courier', fontSize=9, leading=12,
            backColor=rl_colors.HexColor('#f5f5f5'),
            borderWidth=0.5, borderColor=rl_colors.HexColor('#e0e0e0'),
            borderPadding=8, spaceAfter=10, spaceBefore=6
        ))
        styles.add(ParagraphStyle(
            'BulletItem', parent=styles['Normal'],
            fontSize=10.5, leading=15, leftIndent=20, spaceAfter=3,
            bulletIndent=8, bulletFontSize=10
        ))
        styles.add(ParagraphStyle(
            'SubBulletItem', parent=styles['Normal'],
            fontSize=10, leading=14, leftIndent=40, spaceAfter=2,
            bulletIndent=28, bulletFontSize=9
        ))

        flowables = []
        blocks = _parse_markdown_blocks(content)

        for idx, block in enumerate(blocks):
            btype = block['type']

            if btype == 'heading':
                level = block['level']
                text = _md_inline_to_html(block['text'])
                if level == 1 and idx == 0:
                    flowables.append(Paragraph(text, styles['DocTitle']))
                    flowables.append(HRFlowable(width="100%", thickness=2, color=rl_colors.HexColor('#1a237e'), spaceAfter=12))
                elif level == 1:
                    flowables.append(Spacer(1, 6))
                    flowables.append(Paragraph(text, styles['H1']))
                elif level == 2:
                    flowables.append(Paragraph(text, styles['H2']))
                else:
                    flowables.append(Paragraph(text, styles['H3']))

            elif btype == 'paragraph':
                text = _md_inline_to_html(block['text'])
                flowables.append(Paragraph(text, styles['BodyCustom']))

            elif btype == 'bullet_list':
                for item in block['items']:
                    text = _md_inline_to_html(item['text'])
                    style = styles['SubBulletItem'] if item.get('level', 0) > 0 else styles['BulletItem']
                    bullet = '\u2022'
                    if item.get('level', 0) > 0:
                        bullet = '\u25e6'
                    flowables.append(Paragraph(f'{bullet}  {text}', style))

            elif btype == 'numbered_list':
                for num, item in enumerate(block['items'], 1):
                    text = _md_inline_to_html(item['text'])
                    style = styles['BulletItem']
                    flowables.append(Paragraph(f'{num}.  {text}', style))

            elif btype == 'table':
                headers = block['headers']
                rows = block['rows']
                table_data = [headers] + rows
                col_count = len(headers)
                avail_width = letter[0] - 1.5 * inch
                col_width = avail_width / max(col_count, 1)

                t = Table(table_data, colWidths=[col_width] * col_count)
                style_cmds = [
                    ('BACKGROUND', (0, 0), (-1, 0), rl_colors.HexColor('#1a237e')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), rl_colors.white),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('FONTSIZE', (0, 1), (-1, -1), 9.5),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                    ('GRID', (0, 0), (-1, -1), 0.5, rl_colors.HexColor('#bdbdbd')),
                    ('TOPPADDING', (0, 0), (-1, -1), 6),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                    ('LEFTPADDING', (0, 0), (-1, -1), 8),
                ]
                for row_idx in range(1, len(table_data)):
                    if row_idx % 2 == 0:
                        style_cmds.append(('BACKGROUND', (0, row_idx), (-1, row_idx), rl_colors.HexColor('#f5f5f5')))
                t.setStyle(TableStyle(style_cmds))
                flowables.append(Spacer(1, 6))
                flowables.append(t)
                flowables.append(Spacer(1, 8))

            elif btype == 'code':
                code_text = block['text'].replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('\n', '<br/>')
                flowables.append(Paragraph(code_text, styles['CodeBlock']))

            elif btype == 'page_break':
                flowables.append(PageBreak())

            elif btype == 'hr':
                flowables.append(Spacer(1, 4))
                flowables.append(HRFlowable(width="100%", thickness=1, color=rl_colors.HexColor('#bdbdbd'), spaceAfter=8))

        if not flowables:
            return "Error: No content blocks parsed from the provided markdown."

        doc.build(flowables)

        # Verify page count if pages was specified
        page_info = ""
        if pages:
            try:
                import fitz
                verify_doc = fitz.open(filename)
                actual_pages = len(verify_doc)
                verify_doc.close()
                if actual_pages != pages:
                    page_info = f" (Note: Document has {actual_pages} pages, {pages} were requested)"
            except Exception:
                pass

        return f"Successfully created PDF: {filename}{page_info}"
    except Exception as e:
        return f"Error creating PDF: {str(e)}"


def create_docx(filename: str, content: str = "", pages: int = None) -> str:
    """Create a professionally styled DOCX from markdown content."""
    try:
        if not content.strip():
            return "Error: content is required to create a DOCX."

        doc = Document()

        # Set default font
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Calibri'
        font.size = Pt(11)

        blocks = _parse_markdown_blocks(content)

        def _add_formatted_paragraph(doc_obj, text, style_name='Normal', bold=False, italic=False):
            """Add a paragraph with inline markdown formatting (bold/italic)."""
            p = doc_obj.add_paragraph(style=style_name)
            # Split by bold and italic markers
            parts = re.split(r'(\*\*\*.+?\*\*\*|\*\*.+?\*\*|\*.+?\*|`.+?`)', text)
            for part in parts:
                if part.startswith('***') and part.endswith('***'):
                    run = p.add_run(part[3:-3])
                    run.bold = True
                    run.italic = True
                elif part.startswith('**') and part.endswith('**'):
                    run = p.add_run(part[2:-2])
                    run.bold = True
                elif part.startswith('*') and part.endswith('*') and len(part) > 2:
                    run = p.add_run(part[1:-1])
                    run.italic = True
                elif part.startswith('`') and part.endswith('`'):
                    run = p.add_run(part[1:-1])
                    run.font.name = 'Consolas'
                    run.font.size = Pt(9.5)
                    run.font.color.rgb = RGBColor(0x4a, 0x14, 0x8c)
                else:
                    run = p.add_run(part)
                if bold:
                    run.bold = True
                if italic:
                    run.italic = True
            return p

        for idx, block in enumerate(blocks):
            btype = block['type']

            if btype == 'heading':
                level = block['level']
                h = doc.add_heading(block['text'], level=min(level, 4))
                if level <= 2:
                    for run in h.runs:
                        run.font.color.rgb = RGBColor(0x1a, 0x23, 0x7e)

            elif btype == 'paragraph':
                _add_formatted_paragraph(doc, block['text'])

            elif btype == 'bullet_list':
                for item in block['items']:
                    level = item.get('level', 0)
                    style_name = 'List Bullet' if level == 0 else 'List Bullet 2'
                    try:
                        _add_formatted_paragraph(doc, item['text'], style_name)
                    except KeyError:
                        p = _add_formatted_paragraph(doc, item['text'])
                        fmt = p.paragraph_format
                        fmt.left_indent = Inches(0.25 + level * 0.25)

            elif btype == 'numbered_list':
                for item in block['items']:
                    level = item.get('level', 0)
                    style_name = 'List Number' if level == 0 else 'List Number 2'
                    try:
                        _add_formatted_paragraph(doc, item['text'], style_name)
                    except KeyError:
                        p = _add_formatted_paragraph(doc, item['text'])
                        fmt = p.paragraph_format
                        fmt.left_indent = Inches(0.25 + level * 0.25)

            elif btype == 'table':
                headers = block['headers']
                rows = block['rows']
                col_count = len(headers)
                table = doc.add_table(rows=1 + len(rows), cols=col_count, style='Table Grid')
                # Header row
                for ci, header in enumerate(headers):
                    cell = table.rows[0].cells[ci]
                    cell.text = header
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            run.bold = True
                            run.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
                            run.font.size = Pt(10)
                    shading = cell._element.get_or_add_tcPr()
                    bg = shading.makeelement(qn('w:shd'), {
                        qn('w:fill'): '1a237e', qn('w:val'): 'clear'
                    })
                    shading.append(bg)
                # Data rows
                for ri, row in enumerate(rows):
                    for ci, val in enumerate(row):
                        if ci < col_count:
                            table.rows[1 + ri].cells[ci].text = val
                doc.add_paragraph()  # spacing after table

            elif btype == 'code':
                p = doc.add_paragraph()
                run = p.add_run(block['text'])
                run.font.name = 'Consolas'
                run.font.size = Pt(9)
                fmt = p.paragraph_format
                fmt.left_indent = Inches(0.3)
                fmt.space_before = Pt(6)
                fmt.space_after = Pt(6)
                # Gray background via shading XML
                pPr = p._element.get_or_add_pPr()
                shd = pPr.makeelement(qn('w:shd'), {
                    qn('w:fill'): 'f5f5f5', qn('w:val'): 'clear'
                })
                pPr.append(shd)

            elif btype == 'page_break':
                doc.add_page_break()

            elif btype == 'hr':
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(4)
                p.paragraph_format.space_after = Pt(4)
                pPr = p._element.get_or_add_pPr()
                pBdr = pPr.makeelement(qn('w:pBdr'), {})
                bottom = pBdr.makeelement(qn('w:bottom'), {
                    qn('w:val'): 'single', qn('w:sz'): '6',
                    qn('w:space'): '1', qn('w:color'): 'bdbdbd'
                })
                pBdr.append(bottom)
                pPr.append(pBdr)

        doc.save(filename)
        return f"Successfully created DOCX: {filename}"
    except Exception as e:
        return f"Error creating DOCX: {str(e)}"


def create_excel(filename: str, content: str = "", data: list = None, title: str = "") -> str:
    """Create a professionally styled XLSX from content or data."""
    try:
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"

        # Style definitions
        header_font = XlFont(name='Calibri', bold=True, color='FFFFFF', size=11)
        header_fill = PatternFill(start_color='1a237e', end_color='1a237e', fill_type='solid')
        alt_fill = PatternFill(start_color='f5f5f5', end_color='f5f5f5', fill_type='solid')
        title_font = XlFont(name='Calibri', bold=True, size=14, color='1a237e')
        thin_border = Border(
            left=Side(style='thin', color='bdbdbd'),
            right=Side(style='thin', color='bdbdbd'),
            top=Side(style='thin', color='bdbdbd'),
            bottom=Side(style='thin', color='bdbdbd')
        )
        cell_alignment = XlAlignment(vertical='center', wrap_text=True)

        # Determine data source
        table_data = None
        if data:
            table_data = data
        elif content.strip():
            # Try to extract a markdown table from content
            blocks = _parse_markdown_blocks(content)
            for block in blocks:
                if block['type'] == 'table':
                    table_data = [block['headers']] + block['rows']
                    break
            # If no table found, try to split content into rows
            if not table_data:
                lines = [l.strip() for l in content.strip().split('\n') if l.strip()]
                if lines:
                    table_data = []
                    for l in lines:
                        if ',' in l or '\t' in l:
                            sep = '\t' if '\t' in l else ','
                            table_data.append([c.strip() for c in l.split(sep)])
                        else:
                            table_data.append([l])

        if not table_data:
            return "Error: No data provided. Pass either 'content' with a markdown table or 'data' as an array of arrays."

        start_row = 1
        # Add title if provided
        if title:
            ws.cell(row=1, column=1, value=title).font = title_font
            col_count = max(len(row) for row in table_data)
            if col_count > 1:
                ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=col_count)
            start_row = 3

        # Write data
        for ri, row in enumerate(table_data):
            for ci, val in enumerate(row):
                cell = ws.cell(row=start_row + ri, column=ci + 1, value=val)
                cell.border = thin_border
                cell.alignment = cell_alignment
                if ri == 0:
                    cell.font = header_font
                    cell.fill = header_fill
                elif ri % 2 == 0:
                    cell.fill = alt_fill

        # Auto-fit column widths
        for ci in range(1, max(len(row) for row in table_data) + 1):
            max_len = 0
            col_letter = get_column_letter(ci)
            for ri in range(len(table_data)):
                cell_val = str(table_data[ri][ci - 1]) if ci - 1 < len(table_data[ri]) else ''
                max_len = max(max_len, len(cell_val))
            ws.column_dimensions[col_letter].width = min(max(max_len + 4, 10), 50)

        wb.save(filename)
        return f"Successfully created XLSX: {filename}"
    except Exception as e:
        return f"Error creating XLSX: {str(e)}"


def create_pptx(filename: str, content: str = "", pages: int = None) -> str:
    """Create a professionally styled PPTX from markdown content.
    Each heading (# or ##) becomes a new slide title, with content below as bullets."""
    try:
        if not content.strip():
            return "Error: content is required to create a PPTX."

        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        blocks = _parse_markdown_blocks(content)

        # Group blocks into slides: each H1/H2 starts a new slide
        slides_data = []
        current_slide = None

        for block in blocks:
            if block['type'] == 'heading' and block['level'] <= 2:
                if current_slide is not None:
                    slides_data.append(current_slide)
                current_slide = {'title': block['text'], 'blocks': [], 'is_title_slide': (block['level'] == 1 and len(slides_data) == 0)}
            else:
                if current_slide is None:
                    current_slide = {'title': '', 'blocks': [], 'is_title_slide': True}
                current_slide['blocks'].append(block)

        if current_slide is not None:
            slides_data.append(current_slide)

        if not slides_data:
            return "Error: No slides could be generated from the content."

        navy = PptxRGBColor(0x1a, 0x23, 0x7e)
        white = PptxRGBColor(0xff, 0xff, 0xff)
        dark_gray = PptxRGBColor(0x33, 0x33, 0x33)
        light_blue = PptxRGBColor(0x3f, 0x51, 0xb5)

        for si, sd in enumerate(slides_data):
            slide_layout = prs.slide_layouts[6]  # Blank layout for full control
            slide = prs.slides.add_slide(slide_layout)

            if sd.get('is_title_slide'):
                # Title slide with dark background
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = navy

                # Title text box
                txBox = slide.shapes.add_textbox(
                    PptxInches(1), PptxInches(2.2), PptxInches(11.333), PptxInches(1.5)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = sd['title']
                p.font.size = PptxPt(40)
                p.font.bold = True
                p.font.color.rgb = white
                p.alignment = PP_ALIGN.CENTER

                # Subtitle from first paragraph block
                if sd['blocks']:
                    first_block = sd['blocks'][0]
                    if first_block['type'] == 'paragraph':
                        subBox = slide.shapes.add_textbox(
                            PptxInches(2), PptxInches(4), PptxInches(9.333), PptxInches(1)
                        )
                        stf = subBox.text_frame
                        stf.word_wrap = True
                        sp = stf.paragraphs[0]
                        sp.text = re.sub(r'\*+', '', first_block['text'])
                        sp.font.size = PptxPt(20)
                        sp.font.color.rgb = PptxRGBColor(0xbb, 0xbb, 0xff)
                        sp.alignment = PP_ALIGN.CENTER
            else:
                # Content slide
                # Title bar area
                title_shape = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.9)
                )
                tf = title_shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = sd['title']
                p.font.size = PptxPt(28)
                p.font.bold = True
                p.font.color.rgb = navy

                # Underline
                line_shape = slide.shapes.add_shape(
                    1, PptxInches(0.5), PptxInches(1.15), PptxInches(12), Emu(0)
                )
                line_shape.line.color.rgb = light_blue
                line_shape.line.width = PptxPt(2)

                # Content area
                content_box = slide.shapes.add_textbox(
                    PptxInches(0.8), PptxInches(1.5), PptxInches(11.5), PptxInches(5.5)
                )
                ctf = content_box.text_frame
                ctf.word_wrap = True

                first_para = True
                for block in sd['blocks']:
                    if block['type'] == 'paragraph':
                        if first_para:
                            p = ctf.paragraphs[0]
                            first_para = False
                        else:
                            p = ctf.add_paragraph()
                        p.text = re.sub(r'\*+', '', block['text'])
                        p.font.size = PptxPt(16)
                        p.font.color.rgb = dark_gray
                        p.space_after = PptxPt(8)

                    elif block['type'] in ('bullet_list', 'numbered_list'):
                        for bi, item in enumerate(block['items']):
                            if first_para:
                                p = ctf.paragraphs[0]
                                first_para = False
                            else:
                                p = ctf.add_paragraph()
                            clean_text = re.sub(r'\*+', '', item['text'])
                            level = item.get('level', 0)
                            prefix = '\u2022 ' if block['type'] == 'bullet_list' else f'{bi + 1}. '
                            if level > 0:
                                prefix = '   \u25e6 ' if block['type'] == 'bullet_list' else f'   {bi + 1}. '
                            p.text = prefix + clean_text
                            p.font.size = PptxPt(15 if level == 0 else 13)
                            p.font.color.rgb = dark_gray
                            p.space_after = PptxPt(4)
                            p.level = level

                    elif block['type'] == 'heading':
                        if first_para:
                            p = ctf.paragraphs[0]
                            first_para = False
                        else:
                            p = ctf.add_paragraph()
                        p.text = block['text']
                        p.font.size = PptxPt(20)
                        p.font.bold = True
                        p.font.color.rgb = light_blue
                        p.space_before = PptxPt(12)
                        p.space_after = PptxPt(6)

        prs.save(filename)
        return f"Successfully created PPTX: {filename}"
    except Exception as e:
        return f"Error creating PPTX: {str(e)}"


# ─── Document Creation Tool Definitions ──────────────────────────────

CREATE_PDF_DEFINITION = ToolDefinition(
    name="create_pdf",
    description=(
        "Create a professionally styled PDF document from markdown content. "
        "You MUST generate the COMPLETE document content in markdown format and pass it as the 'content' parameter. "
        "Use full markdown: # headings, ## subheadings, **bold**, *italic*, bullet lists (- item), "
        "numbered lists (1. item), tables (| col1 | col2 |), code blocks (```), and --- for horizontal rules. "
        "Write thorough, detailed content - the more content you provide, the better the document.\n\n"
        "PAGE COUNT CONTROL: If the user requests a specific number of pages, you MUST set the 'pages' parameter "
        "AND insert <!-- PAGE_BREAK --> markers in your content to separate each page. "
        "For N pages, include exactly N-1 page break markers. Each section between markers becomes one page. "
        "Write approximately 300-400 words of content per page to fill each page appropriately. "
        "Example for 3 pages: content for page 1 <!-- PAGE_BREAK --> content for page 2 <!-- PAGE_BREAK --> content for page 3"
    ),
    parameters={
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "The path and filename for the PDF (e.g. 'report.pdf')."
            },
            "content": {
                "type": "string",
                "description": "The FULL document content in markdown format. Include headings, paragraphs, lists, tables, etc. Use <!-- PAGE_BREAK --> to force page breaks."
            },
            "pages": {
                "type": "integer",
                "description": "The desired number of pages. When set, you MUST include exactly (pages - 1) <!-- PAGE_BREAK --> markers in the content to control pagination."
            }
        },
        "required": ["filename", "content"]
    },
    function=lambda args: create_pdf(args.get('filename'), args.get('content', ''), args.get('pages'))
)

CREATE_DOCX_DEFINITION = ToolDefinition(
    name="create_docx",
    description=(
        "Create a professionally styled Word document (DOCX) from markdown content. "
        "You MUST generate the COMPLETE document content in markdown format and pass it as the 'content' parameter. "
        "Use full markdown: # headings, ## subheadings, **bold**, *italic*, bullet lists (- item), "
        "numbered lists (1. item), tables (| col1 | col2 |), code blocks (```), and --- for horizontal rules. "
        "Write thorough, detailed content - the more content you provide, the better the document.\n\n"
        "PAGE COUNT CONTROL: If the user requests a specific number of pages, you MUST set the 'pages' parameter "
        "AND insert <!-- PAGE_BREAK --> markers in your content to separate each page. "
        "For N pages, include exactly N-1 page break markers. Each section between markers becomes one page. "
        "Write approximately 350-450 words of content per page to fill each page appropriately. "
        "Example for 3 pages: content for page 1 <!-- PAGE_BREAK --> content for page 2 <!-- PAGE_BREAK --> content for page 3"
    ),
    parameters={
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "The path and filename for the DOCX (e.g. 'report.docx')."
            },
            "content": {
                "type": "string",
                "description": "The FULL document content in markdown format. Include headings, paragraphs, lists, tables, etc. Use <!-- PAGE_BREAK --> to force page breaks."
            },
            "pages": {
                "type": "integer",
                "description": "The desired number of pages. When set, you MUST include exactly (pages - 1) <!-- PAGE_BREAK --> markers in the content to control pagination."
            }
        },
        "required": ["filename", "content"]
    },
    function=lambda args: create_docx(args.get('filename'), args.get('content', ''), args.get('pages'))
)

CREATE_EXCEL_DEFINITION = ToolDefinition(
    name="create_excel",
    description=(
        "Create a professionally styled Excel spreadsheet (XLSX). Provide data either as a markdown table in 'content' "
        "or as an array of arrays in 'data'. The first row is treated as the header row with special styling. "
        "Include an optional 'title' for a merged title row above the data."
    ),
    parameters={
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "The path and filename for the XLSX (e.g. 'data.xlsx')."
            },
            "content": {
                "type": "string",
                "description": "A markdown table to convert into the spreadsheet. e.g. '| Name | Age |\\n|---|---|\\n| Alice | 30 |'"
            },
            "data": {
                "type": "array",
                "items": {"type": "array", "items": {"type": "string"}},
                "description": "Array of arrays representing rows. First row is headers. e.g. [['Name','Age'],['Alice','30']]"
            },
            "title": {
                "type": "string",
                "description": "Optional title displayed as a merged row above the data."
            }
        },
        "required": ["filename"]
    },
    function=lambda args: create_excel(args.get('filename'), args.get('content', ''), args.get('data'), args.get('title', ''))
)

CREATE_PPTX_DEFINITION = ToolDefinition(
    name="create_pptx",
    description=(
        "Create a professionally styled PowerPoint presentation (PPTX) from markdown content. "
        "You MUST generate the COMPLETE presentation content in markdown format. "
        "Each # or ## heading starts a NEW SLIDE with that heading as the slide title. "
        "The first # heading becomes the title slide. Content under each heading becomes bullet points on that slide. "
        "Use bullet lists (- item), numbered lists (1. item), and paragraphs for slide content. "
        "Write concise but informative bullet points for each slide.\n\n"
        "SLIDE COUNT CONTROL: If the user requests a specific number of slides/pages, you MUST set the 'pages' parameter "
        "AND ensure your content contains EXACTLY that many # or ## headings (including the title slide). "
        "For example, for 5 slides: use 1 title heading (# Title) + 4 content headings (## Slide Title). "
        "Count your headings carefully to match the requested number exactly."
    ),
    parameters={
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "The path and filename for the PPTX (e.g. 'presentation.pptx')."
            },
            "content": {
                "type": "string",
                "description": "The FULL presentation content in markdown format. Each # or ## heading starts a new slide."
            },
            "pages": {
                "type": "integer",
                "description": "The desired number of slides. You MUST include EXACTLY this many # or ## headings in your content."
            }
        },
        "required": ["filename", "content"]
    },
    function=lambda args: create_pptx(args.get('filename'), args.get('content', ''), args.get('pages'))
)
